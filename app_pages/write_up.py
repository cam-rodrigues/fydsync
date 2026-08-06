import re
import streamlit as st
import pdfplumber
from calendar import month_name
import pandas as pd
from rapidfuzz import fuzz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from io import BytesIO
import yfinance as yf

#───Performance Table──────────────────────────────────────────────────────────────────

def extract_performance_table(pdf, performance_page, fund_names, end_page=None):
    import re
    from rapidfuzz import fuzz

    # Decide where to stop in the PDF
    end = end_page if end_page is not None else len(pdf.pages) + 1

    # 1. Get all lines from the section
    lines = []
    for pnum in range(performance_page - 1, end - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]

    # 2. Prepare regex
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")

    # 3. For each fund, try to pull numbers
    perf_data = []
    for name in fund_names:
        item = {"Fund Scorecard Name": name}
        tk = ""  # You’ll fill in ticker later from st.session_state["tickers"] or similar
        # a) Exact-ticker match: you’ll want to match this if you have tickers already
        idx = next(
            (i for i, ln in enumerate(lines)
             if name in ln),
            None
        )
        # b) Fuzzy-name fallback if not found
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower()))
                      for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                continue  # Can't find a match

        # c) Pull fund numbers from line above (and two above if needed)
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        clean += [None] * (8 - len(clean))  # pad

        # d) Map to columns
        item["QTD"] = clean[0]
        item["1Yr"] = clean[2]
        item["3Yr"] = clean[3]
        item["5Yr"] = clean[4]
        item["10Yr"] = clean[5]
        item["Net Expense Ratio"] = clean[-2]

        # e) Pull benchmark QTD, 3Yr, 5Yr from next lines
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 1 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]
        item["Bench QTD"] = bench_clean[0] if bench_clean else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None

        perf_data.append(item)
    return perf_data

#───Utility──────────────────────────────────────────────────────────────────

def extract_report_date(text):
    """
    Extracts and formats the report date from a block of text.
    Returns a string like '2nd QTR, 2024' for quarter-end dates,
    or 'As of March 12, 2024' for other dates.
    """
    dates = re.findall(r'(\d{1,2})/(\d{1,2})/(20\d{2})', text or "")
    for month, day, year in dates:
        m, d = int(month), int(day)
        # Quarter-end mapping
        quarter_map = {(3,31): "1st", (6,30): "2nd", (9,30): "3rd", (12,31): "4th"}
        if (m, d) in quarter_map:
            return f"{quarter_map[(m, d)]} QTR, {year}"
        # Fallback: month-day-year as human readable
        return f"As of {month_name[m]} {d}, {year}"
    return None

#───Page 1──────────────────────────────────────────────────────────────────

def process_page1(text):
    """
    Extracts 'report_date', 'total_options', 'prepared_for', and 'prepared_by' from text.
    Populates st.session_state with those keys.
    """
    # Extract report date
    report_date = extract_report_date(text)
    if report_date:
        st.session_state['report_date'] = report_date
    else:
        st.session_state['report_date'] = None

    # Extract total options
    m = re.search(r"Total Options:\s*(\d+)", text or "")
    st.session_state['total_options'] = int(m.group(1)) if m else None

    # Extract prepared for
    m = re.search(r"Prepared For:\s*\n(.*)", text or "")
    st.session_state['prepared_for'] = m.group(1).strip() if m else None

    # Extract prepared by, default if blank or says "mpi stylus"
    m = re.search(r"Prepared By:\s*(.*)", text or "")
    pb = m.group(1).strip() if m else ""
    if not pb or "mpi stylus" in pb.lower():
        pb = "Procyon Partners, LLC"
    st.session_state['prepared_by'] = pb

#───Info Card──────────────────────────────────────────────────────────────────

def show_report_summary():
    report_date    = st.session_state.get('report_date', 'N/A')
    total_options  = st.session_state.get('total_options', 'N/A')
    prepared_for   = st.session_state.get('prepared_for', 'N/A')
    prepared_by    = st.session_state.get('prepared_by', 'N/A')

    st.markdown(f"""
        <div style="
            background: linear-gradient(120deg, #e6f0fb 80%, #c8e0f6 100%);
            color: #244369;
            border-radius: 1.5rem;
            box-shadow: 0 4px 24px rgba(44,85,130,0.11), 0 2px 8px rgba(36,67,105,0.09);
            padding: 1.2rem 2.3rem 1.2rem 2.3rem;
            margin-bottom: 2rem;
            font-size: 1.08rem;
            border: 1.2px solid #b5d0eb;">
            <div style="color:#244369;">
                <b>Report Date:</b> {report_date} <br>
                <b>Total Options:</b> {total_options} <br>
                <b>Prepared For:</b> {prepared_for} <br>
                <b>Prepared By:</b> {prepared_by}
            </div>
        </div>
    """, unsafe_allow_html=True)

#───Step 2: Table of Contents Extraction──────────────────────────────────────────────────────────

def process_toc(text):
    perf = re.search(r"Fund Performance[^\d]*(\d{1,3})", text or "")
    cy   = re.search(r"Fund Performance: Calendar Year\s+(\d{1,3})", text or "")
    r3yr = re.search(r"Risk Analysis: MPT Statistics \(3Yr\)\s+(\d{1,3})", text or "")
    r5yr = re.search(r"Risk Analysis: MPT Statistics \(5Yr\)\s+(\d{1,3})", text or "")

    sc            = re.search(r"Fund Scorecard\s+(\d{1,3})", text or "")
    sc_prop       = re.search(r"Fund Scorecard:\s*Proposed Funds\s+(\d{1,3})", text or "")

    fs            = re.search(r"Fund Factsheets\s+(\d{1,3})", text or "")
    fs_prop       = re.search(r"Fund Factsheets:\s*Proposed Funds\s+(\d{1,3})", text or "")

    perf_page     = int(perf.group(1)) if perf else None
    cy_page       = int(cy.group(1)) if cy else None
    r3yr_page     = int(r3yr.group(1)) if r3yr else None
    r5yr_page     = int(r5yr.group(1)) if r5yr else None
    sc_page       = int(sc.group(1)) if sc else None
    sc_prop_page  = int(sc_prop.group(1)) if sc_prop else None
    fs_page       = int(fs.group(1)) if fs else None
    fs_prop_page  = int(fs_prop.group(1)) if fs_prop else None

    # Store in session state for downstream use
    st.session_state['performance_page'] = perf_page
    st.session_state['calendar_year_page'] = cy_page
    st.session_state['r3yr_page'] = r3yr_page
    st.session_state['r5yr_page'] = r5yr_page
    st.session_state['scorecard_page'] = sc_page
    st.session_state['scorecard_proposed_page'] = sc_prop_page
    st.session_state['factsheets_page'] = fs_page
    st.session_state['factsheets_proposed_page'] = fs_prop_page



#───IPS Invesment Screening──────────────────────────────────────────────────────────────────
import streamlit as st
import re
import pdfplumber
import pandas as pd
import yfinance as yf

def infer_fund_type_guess(ticker):
    """Infer 'Active' or 'Passive' based on Yahoo Finance info (name and summary)."""
    try:
        if not ticker:
            return ""
        info = yf.Ticker(ticker).info
        name = (info.get("longName") or info.get("shortName") or "").lower()
        summary = (info.get("longBusinessSummary") or "").lower()
        if "index" in name or "index" in summary:
            return "Passive"
        if "track" in summary and "index" in summary:
            return "Passive"
        if "actively managed" in summary or "actively-managed" in summary:
            return "Active"
        if "outperform" in summary or "manager selects" in summary:
            return "Active"
        return ""
    except Exception:
        return ""

def extract_scorecard_blocks(pdf, scorecard_page):
    metric_labels = [
        "Manager Tenure", "Excess Performance (3Yr)", "Excess Performance (5Yr)",
        "Peer Return Rank (3Yr)", "Peer Return Rank (5Yr)", "Style Drift Score (3Yr)",
        "Expense Ratio Rank", "Sharpe Ratio Rank (3Yr)", "Sharpe Ratio Rank (5Yr)",
        "R-Squared (3Yr)", "R-Squared (5Yr)",
        "Sortino Ratio Rank (3Yr)", "Sortino Ratio Rank (5Yr)",
        "Tracking Error (3Yr)", "Tracking Error (5Yr)"
    ]
    pages, fund_blocks, fund_name, metrics = [], [], None, []
    # collect all text from scorecard pages
    for p in pdf.pages[scorecard_page-1:]:
        pages.append(p.extract_text() or "")
    lines = "\n".join(pages).splitlines()

    for line in lines:
        # whenever we hit a line that is NOT metric-label-y, it's (potentially) a new fund header
        if not any(metric in line for metric in metric_labels) and line.strip():
            # flush previous fund
            if fund_name and metrics:
                fund_blocks.append({"Fund Name": fund_name, "Metrics": metrics})
            # strip *any* watchlist suffix (14 or 15 criteria, or simple "Meets Watchlist")
            fund_name = re.sub(
                r"Fund\s+(?:Meets Watchlist Criteria|has been placed on watchlist for not meeting .*? out of \d+ criteria)\.?",
                "",
                line.strip(),
                flags=re.IGNORECASE
            ).strip()
            metrics = []

        # now collect any metrics on this line
        for metric in metric_labels:
            if metric in line:
                m = re.match(r"^(.*?)\s+(Pass|Review|Fail)\s*(.*)", line.strip())
                if m:
                    metric_name, status, info = m.groups()
                    metrics.append({
                        "Metric": metric_name,
                        "Status": status,
                        "Info": info.strip()
                    })

    # after loop, append last fund
    if fund_name and metrics:
        fund_blocks.append({"Fund Name": fund_name, "Metrics": metrics})

    return fund_blocks




def extract_fund_tickers(pdf, performance_page, fund_names, factsheets_page=None):
    import re
    from rapidfuzz import fuzz

    def normalize_name(name):
        # strip the watchlist suffix and punctuation, lowercase
        cleaned = re.sub(
            r"has been placed on watchlist for not meeting .*? criteria",
            "",
            name,
            flags=re.IGNORECASE
        )
        cleaned = re.sub(r"[^A-Za-z0-9 ]+", "", cleaned)  # remove punctuation
        return cleaned.strip().lower()

    end_page = factsheets_page - 1 if factsheets_page else len(pdf.pages)
    all_lines = []
    for p in pdf.pages[performance_page - 1:end_page]:
        txt = p.extract_text() or ""
        all_lines.extend([ln.strip() for ln in txt.splitlines() if ln.strip()])

    # Step 1: Collect candidate (raw_name, ticker) pairs from lines with uppercase tickers length 2-5
    candidate_pairs = []  # list of tuples (normalized_raw_name, ticker, raw_name_original)
    ticker_rx = re.compile(r"\b([A-Z]{2,5})\b")
    for ln in all_lines:
        matches = ticker_rx.findall(ln)
        if not matches:
            continue
        for ticker in set(matches):
            parts = ln.rsplit(ticker, 1)
            if len(parts) >= 1:
                raw_name = parts[0].strip()
                if not raw_name:
                    continue
                norm_raw = normalize_name(raw_name)
                if not norm_raw:
                    continue
                candidate_pairs.append((norm_raw, ticker, raw_name))

    # Step 2: For each fund_name, find best candidate fuzzy match (one-to-one)
    assigned = {}
    used_tickers = set()

    # Precompute normalized expected names
    norm_expected_list = [
        (name, normalize_name(name))
        for name in fund_names
    ]

    for fund_name, norm_expected in norm_expected_list:
        best = (None, None, 0)  # (ticker, raw_name_original, score)
        for norm_raw, ticker, raw_name in candidate_pairs:
            if ticker in used_tickers:
                continue
            score = fuzz.token_sort_ratio(norm_expected, norm_raw)
            if score > best[2]:
                best = (ticker, raw_name, score)
        if best[2] >= 70:  # threshold, adjust if needed
            assigned[fund_name] = best[0]
            used_tickers.add(best[0])
        else:
            assigned[fund_name] = ""

    # Step 3: Fallback for unmatched names: looser heuristic and line-based fallback
    for name in fund_names:
        if assigned.get(name):
            continue
        norm_expected = normalize_name(name)
        for norm_raw, ticker, raw_name in candidate_pairs:
            if ticker in used_tickers:
                continue
            if norm_expected in norm_raw or norm_raw in norm_expected:
                assigned[name] = ticker
                used_tickers.add(ticker)
                break
        if not assigned.get(name):
            for ln in all_lines:
                if name.lower() in ln.lower():
                    m = re.search(r"\b([A-Z]{2,5})\b", ln)
                    if m:
                        assigned[name] = m.group(1)
                        break

    # Final cleanup: ensure non-None strings
    return {k: (v if v else "") for k, v in assigned.items()}

def scorecard_to_ips(fund_blocks, fund_types, tickers):
    """
    Converts raw scorecard metrics into the 11 IPS checks for each fund,
    working correctly whether the scorecard has 14 or 15 total metrics.
    """

    # Define, in order, which scorecard metric each IPS check corresponds to.
    ACTIVE_IPS_METRICS = [
        "Manager Tenure",
        "Excess Performance (3Yr)",
        "Peer Return Rank (3Yr)",
        "Sharpe Ratio Rank (3Yr)",
        "Sortino Ratio Rank (3Yr)",
        "Excess Performance (5Yr)",
        "Peer Return Rank (5Yr)",
        "Sharpe Ratio Rank (5Yr)",
        "Sortino Ratio Rank (5Yr)",
        "Expense Ratio Rank",
        # last check is always a Pass
        None
    ]
    PASSIVE_IPS_METRICS = [
        "Manager Tenure",
        "R-Squared (3Yr)",
        "Peer Return Rank (3Yr)",
        "Sharpe Ratio Rank (3Yr)",
        "Tracking Error (3Yr)",
        "R-Squared (5Yr)",
        "Peer Return Rank (5Yr)",
        "Sharpe Ratio Rank (5Yr)",
        "Tracking Error (5Yr)",
        "Expense Ratio Rank",
        None
    ]

    ips_labels = [f"IPS Investment Criteria {i+1}" for i in range(11)]
    icon_rows, raw_rows = [], []

    for fund in fund_blocks:
        name      = fund["Fund Name"]
        ftype     = fund_types.get(name, 
                       "Passive" if "index" in name.lower() else "Active"
                   )
        status_map = {m["Metric"]: m["Status"] for m in fund["Metrics"]}

        # pick the right mapping by name
        metric_names = PASSIVE_IPS_METRICS if ftype == "Passive" else ACTIVE_IPS_METRICS

        # build the 11 statuses in order
        ips_status = []
        for metric_name in metric_names:
            if metric_name is None:
                ips_status.append("Pass")
            else:
                ips_status.append(status_map.get(metric_name, "Pass"))

        # count failures
        fails = sum(1 for s in ips_status if s in ("Review", "Fail"))
        watch = "FW" if fails >= 6 else "IW" if fails >= 5 else "NW"

        # map to icons
        def ico(s):
            return "✔" if s == "Pass" else ("✗" if s in ("Review", "Fail") else "")

        icon_row = {
            "Fund Name":       name,
            "Ticker":          tickers.get(name, ""),
            "Fund Type":       ftype,
            **{ips_labels[i]: ico(ips_status[i]) for i in range(11)},
            "IPS Watch Status": watch
        }
        raw_row = {
            "Fund Name":       name,
            "Ticker":          tickers.get(name, ""),
            "Fund Type":       ftype,
            **{ips_labels[i]: ips_status[i] for i in range(11)},
            "IPS Watch Status": watch
        }

        icon_rows.append(icon_row)
        raw_rows.append(raw_row)

    return pd.DataFrame(icon_rows), pd.DataFrame(raw_rows)



def watch_status_color(val):
    if val == "FW":
        return "background-color:#f8d7da; color:#c30000; font-weight:600;"
    if val == "IW":
        return "background-color:#fff3cd; color:#B87333; font-weight:600;"
    if val == "NW":
        return "background-color:#d6f5df; color:#217a3e; font-weight:600;"
    return ""

def step3_5_6_scorecard_and_ips(pdf, scorecard_page, performance_page, factsheets_page, total_options):
    import pandas as pd
    import streamlit as st

    # --- 1. Extract scorecard blocks ---
    fund_blocks = extract_scorecard_blocks(pdf, scorecard_page)
    fund_names = [fund["Fund Name"] for fund in fund_blocks]
    if not fund_blocks:
        st.error("Could not extract fund scorecard blocks. Check the PDF and page number.")
        return

    # --- 2. Extract tickers ---
    tickers = extract_fund_tickers(pdf, performance_page, fund_names, factsheets_page)

    # --- Prepare inferred fund type guesses/defaults ---
    inferred_guesses = []
    for name in fund_names:
        guess = ""
        if tickers.get(name):
            guess = infer_fund_type_guess(tickers.get(name, "")) or ""
        inferred_guesses.append("Passive" if guess.lower() == "passive" else ("Passive" if "index" in name.lower() else "Active"))

    # Build base df for editor
    df_types_base = pd.DataFrame({
        "Fund Name":       fund_names,
        "Ticker":          [tickers.get(n, "") for n in fund_names],
        "Inferred Type":   inferred_guesses,
        "Fund Type":       inferred_guesses,  # starts same as inferred
    })

    # --- Toggleable editor display ---
    if "show_edit_fund_type" not in st.session_state:
        st.session_state["show_edit_fund_type"] = False
    toggle_label = "Hide Fund Type Editor" if st.session_state["show_edit_fund_type"] else "Edit Fund Type"
    if st.button(toggle_label, key="toggle_edit_fund_type"):
        st.session_state["show_edit_fund_type"] = not st.session_state["show_edit_fund_type"]

    # Decide fund_types mapping
    if st.session_state["show_edit_fund_type"]:
        st.markdown("### Fund Type Overrides")
        st.caption("Edit any fund type here; once you modify at least one, your edits take precedence over inferred types.")
        edited_types = st.data_editor(
            df_types_base,
            column_config={
                "Fund Type": st.column_config.SelectboxColumn("Fund Type", options=["Active", "Passive"]),
            },
            disabled=["Fund Name", "Ticker", "Inferred Type"],
            hide_index=True,
            key="data_editor_fundtype_ips",
            use_container_width=True,
        )
        # Determine whether any manual edit occurred (diff between Fund Type and Inferred Type)
        manual_override = any(
            row["Fund Type"] != row["Inferred Type"]
            for _, row in edited_types.iterrows()
        )
        if manual_override:
            fund_types = {row["Fund Name"]: row["Fund Type"] for _, row in edited_types.iterrows()}
        else:
            # none edited: use inferred
            fund_types = {row["Fund Name"]: row["Inferred Type"] for _, row in edited_types.iterrows()}
    else:
        # editor hidden: use inferred guesses
        fund_types = {name: inferred_guesses[i] for i, name in enumerate(fund_names)}

    # --- 4. IPS conversion ---
    df_icon, df_raw = scorecard_to_ips(fund_blocks, fund_types, tickers)

    # --- IPS Results ---
    st.subheader("IPS Screening Results")
    st.markdown(
        '<div style="display:flex; gap:1rem; margin-bottom:0.5rem;">'
        '<div style="padding:4px 10px; background:#d6f5df; border-radius:4px; font-weight:600;">NW: No Watch</div>'
        '<div style="padding:4px 10px; background:#fff3cd; border-radius:4px; font-weight:600;">IW: Informal Watch</div>'
        '<div style="padding:4px 10px; background:#f8d7da; border-radius:4px; font-weight:600;">FW: Formal Watch</div>'
        '</div>', unsafe_allow_html=True
    )

    if df_icon.empty:
        st.info("No IPS screening data available.")
    else:
        # Compact numbered criteria
        display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
        display_df = df_icon.rename(columns=display_columns)

        def iconify(s):
            if s == "Pass":
                return "✔"
            if s in ("Review", "Fail"):
                return "✗"
            return ""

        compact_df = display_df.copy()
        for i in range(1, 12):
            orig = f"IPS Investment Criteria {i}"
            if orig in compact_df.columns:
                compact_df[str(i)] = compact_df[orig].apply(iconify)
                compact_df.drop(columns=[orig], inplace=True)
        cols_order = ["Fund Name", "Ticker", "Fund Type"] + [str(i) for i in range(1, 12)] + ["IPS Watch Status"]
        compact_df = compact_df[[c for c in cols_order if c in compact_df.columns]]

        def watch_style(val):
            if val == "NW":
                return "background-color:#d6f5df; color:#217a3e; font-weight:600;"
            if val == "IW":
                return "background-color:#fff3cd; color:#B87333; font-weight:600;"
            if val == "FW":
                return "background-color:#f8d7da; color:#c30000; font-weight:600;"
            return ""

        styled = compact_df.style.applymap(watch_style, subset=["IPS Watch Status"])
        st.dataframe(styled, use_container_width=True, hide_index=True)

        # --- Summary badges (bottom) ---
        def summarize_watch(df):
            counts = df["IPS Watch Status"].value_counts().to_dict()
            return {
                "No Watch": counts.get("NW", 0),
                "Informal Watch": counts.get("IW", 0),
                "Formal Watch": counts.get("FW", 0),
            }

        summary = summarize_watch(df_icon)
        st.markdown("---")
        st.markdown("### Watch Summary")
        b1, b2, b3 = st.columns(3, gap="small")
        with b1:
            st.metric("No Watch", summary["No Watch"])
        with b2:
            st.metric("Informal Watch", summary["Informal Watch"])
        with b3:
            st.metric("Formal Watch", summary["Formal Watch"])

    # --- Persist state downstream ---
    st.session_state["fund_blocks"] = fund_blocks
    st.session_state["fund_types"] = fund_types
    st.session_state["fund_tickers"] = tickers
    st.session_state["ips_icon_table"] = df_icon
    st.session_state["ips_raw_table"] = df_raw

    # --- Performance extraction ---
    if performance_page is None:
        st.error("❌ ‘Fund Performance’ page number not found in TOC. Run Step 2 first.")
        perf_data = []
    else:
        perf_data = extract_performance_table(
            pdf,
            performance_page,
            fund_names,
            factsheets_page
        )

    for itm in perf_data:
        itm["Ticker"] = tickers.get(itm["Fund Scorecard Name"], "")
    st.session_state["fund_performance_data"] = perf_data

    st.session_state["tickers"] = tickers  # legacy compatibility


#───Step 6:Factsheets Pages──────────────────────────────────────────────────────────────────

def step6_process_factsheets(pdf, fund_names, suppress_output=True):
    # If you ever want UI, set suppress_output=False when calling

    factsheet_start = st.session_state.get("factsheets_page")
    total_declared = st.session_state.get("total_options")
    performance_data = [
        {"Fund Scorecard Name": name, "Ticker": ticker}
        for name, ticker in st.session_state.get("tickers", {}).items()
    ]

    if not factsheet_start:
        if not suppress_output:
            st.error("❌ 'Fund Factsheets' page number not found in TOC.")
        return

    matched_factsheets = []
    for i in range(factsheet_start - 1, len(pdf.pages)):
        page = pdf.pages[i]
        words = page.extract_words(use_text_flow=True)
        header_words = [w['text'] for w in words if w['top'] < 100]
        first_line = " ".join(header_words).strip()

        if not first_line or "Benchmark:" not in first_line or "Expense Ratio:" not in first_line:
            continue

        ticker_match = re.search(r"\b([A-Z]{5})\b", first_line)
        ticker = ticker_match.group(1) if ticker_match else ""
        fund_name_raw = first_line.split(ticker)[0].strip() if ticker else first_line

        best_score = 0
        matched_name = matched_ticker = ""
        for item in performance_data:
            ref = f"{item['Fund Scorecard Name']} {item['Ticker']}".strip()
            score = fuzz.token_sort_ratio(f"{fund_name_raw} {ticker}".lower(), ref.lower())
            if score > best_score:
                best_score, matched_name, matched_ticker = score, item['Fund Scorecard Name'], item['Ticker']

        def extract_field(label, text, stop=None):
            try:
                start = text.index(label) + len(label)
                rest = text[start:]
                if stop and stop in rest:
                    return rest[:rest.index(stop)].strip()
                return rest.split()[0]
            except Exception:
                return ""

        benchmark = extract_field("Benchmark:", first_line, "Category:")
        category  = extract_field("Category:", first_line, "Net Assets:")
        net_assets= extract_field("Net Assets:", first_line, "Manager Name:")
        manager   = extract_field("Manager Name:", first_line, "Avg. Market Cap:")
        avg_cap   = extract_field("Avg. Market Cap:", first_line, "Expense Ratio:")
        expense   = extract_field("Expense Ratio:", first_line)

        matched_factsheets.append({
            "Page #": i + 1,
            "Parsed Fund Name": fund_name_raw,
            "Parsed Ticker": ticker,
            "Matched Fund Name": matched_name,
            "Matched Ticker": matched_ticker,
            "Benchmark": benchmark,
            "Category": category,
            "Net Assets": net_assets,
            "Manager Name": manager,
            "Avg. Market Cap": avg_cap,
            "Expense Ratio": expense,
            "Match Score": best_score,
            "Matched": "✅" if best_score > 20 else "❌"
        })

    df_facts = pd.DataFrame(matched_factsheets)
    st.session_state['fund_factsheets_data'] = matched_factsheets

    # Hide UI output unless suppress_output is False
    if not suppress_output:
        display_df = df_facts[[
            "Matched Fund Name", "Matched Ticker", "Benchmark", "Category",
            "Net Assets", "Manager Name", "Avg. Market Cap", "Expense Ratio", "Matched"
        ]].rename(columns={"Matched Fund Name": "Fund Name", "Matched Ticker": "Ticker"})

        st.dataframe(display_df, use_container_width=True)

        matched_count = sum(1 for r in matched_factsheets if r["Matched"] == "✅")
        st.write(f"Matched {matched_count} of {len(matched_factsheets)} factsheet pages.")
        if matched_count == total_declared:
            st.success(f"All {matched_count} funds matched the declared Total Options from Page 1.")
        else:
            st.error(f"Mismatch: Page 1 declared {total_declared}, but only matched {matched_count}.")

#───Step 7: QTD / 1Yr / 3Yr / 5Yr / 10Yr / Net Expense Ratio & Bench QTD──────────────────────────────────────────────────────────────────

def step7_extract_returns(pdf):
    import re
    import pandas as pd
    import streamlit as st
    from rapidfuzz import fuzz

    # 1) Where to scan
    perf_page = st.session_state.get("performance_page")
    end_page  = st.session_state.get("calendar_year_page") or (len(pdf.pages) + 1)
    perf_data = st.session_state.get("fund_performance_data", [])
    if perf_page is None or not perf_data:
        st.error("❌ Run Step 5 first to populate performance data.")
        return

    # 2) Prep output slots
    fields = [
        "QTD", "1Yr", "3Yr", "5Yr", "10Yr", "Net Expense Ratio",
        "Bench QTD", "Bench 1Yr", "Bench 3Yr", "Bench 5Yr", "Bench 10Yr"
    ]
    for itm in perf_data:
        for f in fields:
            itm.setdefault(f, None)

    # 3) Gather every nonblank line in the Performance section
    lines = []
    for pnum in range(perf_page - 1, end_page - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]

    # 4) Regex to pull decimal tokens (with optional % and parentheses)
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")

    matched = 0
    for item in perf_data:
        name = item["Fund Scorecard Name"]
        tk   = item["Ticker"].upper().strip()

        # a) Exact-ticker match
        idx = next(
            (i for i, ln in enumerate(lines)
             if re.search(rf"\b{re.escape(tk)}\b", ln)),
            None
        )
        # b) Fuzzy-name fallback
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower()))
                      for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                st.warning(f"⚠️ {name} ({tk}): no match found.")
                continue

        # c) Pull fund numbers from line above (and two above if needed)
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        if len(clean) < 8:
            clean += [None] * (8 - len(clean))

        # d) Map fund returns & net expense
        item["QTD"]               = clean[0]
        item["1Yr"]               = clean[2]
        item["3Yr"]               = clean[3]
        item["5Yr"]               = clean[4]
        item["10Yr"]              = clean[5]
        item["Net Expense Ratio"] = clean[-2]

        # e) Pull benchmark QTD, 1Yr, 3Yr, 5Yr, and 10Yr from the very next line(s)
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 5 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]

        item["Bench QTD"]  = bench_clean[0] if len(bench_clean) > 0 else None
        item["Bench 1Yr"] = bench_clean[1] if len(bench_clean) > 1 else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None
        item["Bench 10Yr"] = bench_clean[5] if len(bench_clean) > 5 else None

        matched += 1

    # 5) Save & display
    st.session_state["fund_performance_data"] = perf_data
    df = pd.DataFrame(perf_data)
    # Hide the per-fund warnings and overall success message, but still alert if something is off
    
    expected_count = len(perf_data)
    if matched < expected_count:
        st.error(f"❌ Only matched {matched} of {expected_count} funds with return data. Check your PDF or extraction logic.")
    # (Do NOT display per-fund warnings or success)

    st.dataframe(
        df[["Fund Scorecard Name", "Ticker"] + fields],
        use_container_width=True
    )



# ─── Step 8: Calendar Year Returns (auto-discover pages) ─────────────────────────────────
def step8_calendar_returns(pdf):
    # 1) Auto-find start page if needed
    cy_page = st.session_state.get("calendar_year_page")
    if cy_page is None:
        for i, pg in enumerate(pdf.pages, start=1):
            if "Fund Performance: Calendar Year" in (pg.extract_text() or ""):
                cy_page = i
                st.session_state["calendar_year_page"] = i
                break
    if cy_page is None:
        st.error("❌ 'Fund Performance: Calendar Year' not found in PDF.")
        return

    # 2) Auto-find end page via 3Yr Risk section
    end_page = st.session_state.get("r3yr_page")
    if end_page is None:
        for i, pg in enumerate(pdf.pages, start=1):
            if "Risk Analysis: MPT Statistics (3Yr)" in (pg.extract_text() or ""):
                end_page = i
                st.session_state["r3yr_page"] = i
                break
    end_page = end_page or (len(pdf.pages) + 1)

    # 3) Extract lines between cy_page and end_page
    lines = []
    for pg in pdf.pages[cy_page-1:end_page-1]:
        txt = pg.extract_text() or ""
        lines.extend([ln.strip() for ln in txt.splitlines() if ln.strip()])

    # 4) Header + years
    header = next((ln for ln in lines if "Ticker" in ln and re.search(r"20\d{2}", ln)), None)
    if not header:
        st.error("❌ Couldn’t find header row with 'Ticker' + year.")
        return
    years = re.findall(r"\b20\d{2}\b", header)
    num_rx = re.compile(r"-?\d+\.\d+%?")

    # A) Funds
    fund_map = st.session_state.get("tickers", {})
    fund_records = []
    for name, tk in fund_map.items():
        ticker = (tk or "").upper().strip()
        idx = next((i for i, ln in enumerate(lines) if ticker and ticker in ln.split()), None)
        if idx is None:
            st.warning(f"⚠️ Calendar returns: no line for {name} ({ticker})")
            continue
        prior = lines[idx-1] if idx >= 1 else ""
        raw   = num_rx.findall(prior)
        vals  = raw[:len(years)] + [None]*max(0, len(years)-len(raw))
        rec   = {"Name": name, "Ticker": ticker}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        fund_records.append(rec)

    df_fund = pd.DataFrame(fund_records)
    if not df_fund.empty:
        st.markdown("**Fund Calendar-Year Returns**")
        st.dataframe(df_fund[["Name","Ticker"]+years], use_container_width=True)
        st.session_state["step8_returns"] = fund_records

    # B) Benchmarks
    bench_records = []
    for f in st.session_state.get("fund_factsheets_data", []):
        bench_name = f.get("Benchmark","").strip()
        bench_tkr  = f.get("Matched Ticker","").upper().strip()
        if not bench_name:
            continue
        idx = next((i for i, ln in enumerate(lines) if bench_name in ln), None)
        if idx is None:
            st.warning(f"⚠️ Calendar returns: no line for benchmark '{bench_name}'")
            continue
        raw  = num_rx.findall(lines[idx])
        vals = raw[:len(years)] + [None]*max(0,len(years)-len(raw))
        rec  = {"Name": bench_name, "Ticker": bench_tkr}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        bench_records.append(rec)

    df_bench = pd.DataFrame(bench_records)
    if not df_bench.empty:
        st.markdown("**Benchmark Calendar-Year Returns**")
        st.dataframe(df_bench[["Name","Ticker"]+years], use_container_width=True)
        st.session_state["benchmark_calendar_year_returns"] = bench_records
    else:
        st.warning("No benchmark returns extracted.")


# ─── Step 9: 3-Yr Risk Analysis (auto-discover start page) ───────────────────────────────
def step9_risk_analysis_3yr(pdf):
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # auto-find 3Yr section
    start = st.session_state.get("r3yr_page")
    if start is None:
        for i, pg in enumerate(pdf.pages, start=1):
            if "Risk Analysis: MPT Statistics (3Yr)" in (pg.extract_text() or ""):
                start = i
                st.session_state["r3yr_page"] = i
                break
    if start is None:
        st.error("❌ ‘Risk Analysis: MPT Statistics (3Yr)’ page not found; run Step 2 first.")
        return

    # scan forward until all tickers seen
    locs = {}
    for pnum in range(start, len(pdf.pages)+1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            toks = ln.split()
            for name, tk in fund_map.items():
                if name in locs: continue
                if tk.upper() in toks:
                    locs[name] = (pnum, li)
        if len(locs)==len(fund_map):
            break

    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, (pgn, li) in locs.items():
        lines = (pdf.pages[pgn-1].extract_text() or "").splitlines()
        nums  = num_rx.findall(lines[li])
        nums += [None]*max(0,4-len(nums))
        alpha,beta,up,down = nums[:4]
        results.append({
            "Fund Name": name,
            "Ticker":    fund_map[name].upper(),
            "3 Year Alpha":          alpha,
            "3 Year Beta":           beta,
            "3 Year Upside Capture": up,
            "3 Year Downside Capture": down,
        })

    st.session_state["step9_mpt_stats"] = results


# ─── Step 10: 5-Yr Risk Analysis (auto-discover start page) ─────────────────────────────
def step10_risk_analysis_5yr(pdf):
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # auto-find 5Yr section
    start = None
    for i, pg in enumerate(pdf.pages, start=1):
        if "Risk Analysis: MPT Statistics (5Yr)" in (pg.extract_text() or ""):
            start = i
            st.session_state["r5yr_page"] = i
            break
    if start is None:
        st.error("❌ Could not find ‘Risk Analysis: MPT Statistics (5Yr)’ section.")
        return

    # scan forward until all tickers seen
    locs = {}
    for pnum in range(start, len(pdf.pages)+1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            toks = ln.split()
            for name, tk in fund_map.items():
                if name in locs: continue
                if tk.upper() in toks:
                    locs[name] = (pnum, li)
        if len(locs)==len(fund_map):
            break

    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, (pgn, li) in locs.items():
        lines = (pdf.pages[pgn-1].extract_text() or "").splitlines()
        nums  = []
        for j in range(li, min(li+3, len(lines))):
            nums += num_rx.findall(lines[j])
            if len(nums)>=4: break
        nums += [None]*max(0,4-len(nums))
        a5,b5,u5,d5 = nums[:4]
        results.append({
            "Fund Name": name,
            "Ticker":    fund_map[name].upper(),
            "5 Year Alpha":          a5,
            "5 Year Beta":           b5,
            "5 Year Upside Capture": u5,
            "5 Year Downside Capture": d5,
        })

    st.session_state["step10_mpt_stats"] = results

#───Step 11: Combined MPT Statistics Summary──────────────────────────────────────────────────────────────────

def step11_create_summary(pdf=None):
    import pandas as pd
    import streamlit as st

    # 1) Load your 3‑Yr and 5‑Yr stats from session state
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    mpt5 = st.session_state.get("step10_mpt_stats", [])
    if not mpt3 or not mpt5:
        st.error("❌ Missing MPT stats. Run Steps 9 & 10 first.")
        return

    # 2) Build DataFrames
    df3 = pd.DataFrame(mpt3)  # contains "3 Year Alpha", "3 Year Beta", etc.
    df5 = pd.DataFrame(mpt5)  # contains "5 Year Alpha", "5 Year Beta", etc.

    # 3) Merge on Fund Name & Ticker
    df = pd.merge(
        df3,
        df5,
        on=["Fund Name", "Ticker"],
        how="outer",
        suffixes=("_3yr", "_5yr")
    )

    # 4) Build the Investment Manager column
    df.insert(0, "Investment Manager", df["Fund Name"] + " (" + df["Ticker"] + ")")

    # 5) Select & order the columns
    df = df[[
        "Investment Manager",
        "3 Year Alpha",
        "5 Year Alpha",
        "3 Year Beta",
        "5 Year Beta",
        "3 Year Upside Capture",
        "3 Year Downside Capture",
        "5 Year Upside Capture",
        "5 Year Downside Capture"
    ]]

    # 6) Display
    st.session_state["step11_summary"] = df.to_dict("records")
    st.dataframe(df)

#───Step 12: Extract “FUND FACTS” & Its Table Details in One Go──────────────────────────────────────────────────────────────────

def step12_process_fund_facts(pdf):
    import re
    import streamlit as st
    import pandas as pd

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # the exact labels and the order you want them in the table
    labels = [
        "Manager Tenure Yrs.",
        "Expense Ratio",
        "Expense Ratio Rank",
        "Total Number of Holdings",
        "Turnover Ratio"
    ]

    records = []
    # scan each factsheet page
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = pdf.pages[pnum-1].extract_text().splitlines()

        for idx, line in enumerate(lines):
            if line.lstrip().upper().startswith("FUND FACTS"):
                # grab the next 8 lines (should contain your 5 labels)
                snippet = lines[idx+1 : idx+1+8]
                rec = {"Fund Name": fund_name, "Ticker": ticker}
                for lab in labels:
                    val = None
                    for ln in snippet:
                        norm = " ".join(ln.strip().split())
                        if norm.startswith(lab):
                            rest = norm[len(lab):].strip(" :\t")
                            m = re.match(r"(-?\d+\.\d+)", rest)
                            val = m.group(1) if m else (rest.split()[0] if rest else None)
                            break
                    rec[lab] = val
                records.append(rec)
                break  # move on to the next page once Fund Facts is processed

    if not records:
        st.warning("No Fund Facts tables found.")
        return

    # save & show
    st.session_state["step12_fund_facts_table"] = records

#───Step 13: Extract Risk‑Adjusted Returns Metrics──────────────────────────────────────────────────────────────────

def step13_process_risk_adjusted_returns(pdf):
    import re
    import streamlit as st
    import pandas as pd

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # which metrics to pull
    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    num_rx  = re.compile(r"-?\d+\.\d+")

    records = []
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()

        # find the heading
        for idx, line in enumerate(lines):
            norm = " ".join(line.strip().split()).upper()
            if norm.startswith("RISK-ADJUSTED RETURNS"):
                snippet = lines[idx+1 : idx+1+6]  # grab next few lines
                rec = {"Fund Name": fund_name, "Ticker": ticker}

                for metric in metrics:
                    # find the snippet line for this metric
                    text_line = next(
                        ( " ".join(ln.strip().split())
                          for ln in snippet
                          if ln.strip().upper().startswith(metric.upper()) ),
                        None
                    ) or ""
                    # extract up to 4 numbers
                    nums = num_rx.findall(text_line)
                    nums += [None] * (4 - len(nums))

                    # assign into rec
                    rec[f"{metric} 1Yr"]  = nums[0]
                    rec[f"{metric} 3Yr"]  = nums[1]
                    rec[f"{metric} 5Yr"]  = nums[2]
                    rec[f"{metric} 10Yr"] = nums[3]

                records.append(rec)
                break  # done with this page

    if not records:
        st.warning("No 'RISK‑ADJUSTED RETURNS' tables found.")
        return

    # save & show
    st.session_state["step13_risk_adjusted_table"] = records
    df = pd.DataFrame(records)
    st.dataframe(df, use_container_width=True)

#───Step 14: Peer Risk-Adjusted Return Rank──────────────────────────────────────────────────────────────────

def step14_extract_peer_risk_adjusted_return_rank(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.write("Peer Rank")

    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    page_map = {
        f["Page #"]:(f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    records = []

    for pnum, (fund, ticker) in page_map.items():
        page = pdf.pages[pnum-1]
        text = page.extract_text() or ""
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]

        # 1) locate the Risk-Adjusted Returns header
        try:
            risk_idx = next(i for i, ln in enumerate(lines)
                            if "RISK-ADJUSTED RETURNS" in ln.upper())
        except StopIteration:
            st.warning(f"⚠️ {fund} ({ticker}): Risk-Adjusted header not found.")
            continue

        # 2) find all “1 Yr 3 Yrs 5 Yrs 10 Yrs” lines after that
        header_idxs = [i for i, ln in enumerate(lines)
                       if re.match(r"1\s*Yr", ln)]
        peer_header_idxs = [i for i in header_idxs if i > risk_idx]

        if not peer_header_idxs:
            st.warning(f"⚠️ {fund} ({ticker}): peer header not found.")
            continue

        # take the *second* header occurrence (first is Risk-Adjusted, next is Peer)
        peer_hdr = peer_header_idxs[0] if len(peer_header_idxs)==1 else peer_header_idxs[1]

        rec = {"Fund Name": fund, "Ticker": ticker}

        # 3) read the three lines immediately below that header
        for offset, metric in enumerate(metrics, start=1):
            if peer_hdr + offset < len(lines):
                parts = lines[peer_hdr + offset].split()
                # parts[0:2] = metric name words, parts[2:6] = the four integer ranks
                vals = parts[2:6] if len(parts) >= 6 else []
            else:
                vals = []

            # fill into record (pad with None if missing)
            for idx, period in enumerate(["1Yr","3Yr","5Yr","10Yr"]):
                rec[f"{metric} {period}"] = vals[idx] if idx < len(vals) else None

            if len(vals) < 4:
                st.warning(f"⚠️ {fund} ({ticker}): only {len(vals)} peer values found for '{metric}'.")

        records.append(rec)

    if not records:
        st.warning("❌ No Peer Risk-Adjusted Return Rank data extracted.")
        return

    df = pd.DataFrame(records)
    st.session_state["step14_peer_rank_table"] = records
    st.dataframe(df, use_container_width=True)

#───Step 14.5: IPS Fail Table──────────────────────────────────────────────────────────────────

def step14_5_ips_fail_table():
    import pandas as pd
    import streamlit as st

    df = st.session_state.get("ips_icon_table")

    if df is None or df.empty:
        return

    fail_df = df[
        df["IPS Watch Status"].isin(["FW", "IW"])
    ][
        [
            "Fund Name",
            "IPS Watch Status",
        ]
    ].copy()

    if fail_df.empty:
        st.success("No funds are currently on watch.")
        return

    status_labels = {
        "IW": "Informal Watch",
        "FW": "Formal Watch",
    }

    fail_df["Watch Status"] = fail_df[
        "IPS Watch Status"
    ].map(status_labels)

    informal_count = int(
        (fail_df["IPS Watch Status"] == "IW").sum()
    )

    formal_count = int(
        (fail_df["IPS Watch Status"] == "FW").sum()
    )

    total_count = len(fail_df)

    st.markdown("### Funds on Watch")
    st.caption(
        "Funds that failed five or more IPS criteria "
        "and require additional review."
    )

    total_col, informal_col, formal_col = st.columns(
        3,
        gap="small",
    )

    with total_col:
        st.metric(
            "Total on Watch",
            total_count,
        )

    with informal_col:
        st.metric(
            "Informal Watch",
            informal_count,
        )

    with formal_col:
        st.metric(
            "Formal Watch",
            formal_count,
        )

    display_df = fail_df.rename(
        columns={
            "Fund Name": "Fund",
        }
    )[
        [
            "Fund",
            "Watch Status",
        ]
    ]

    def style_watch_status(value):
        if value == "Formal Watch":
            return (
                "background-color: #FDECEC; "
                "color: #B42318; "
                "font-weight: 700;"
            )

        if value == "Informal Watch":
            return (
                "background-color: #FFF6DB; "
                "color: #8A6100; "
                "font-weight: 700;"
            )

        return ""

    styled_df = display_df.style.applymap(
        style_watch_status,
        subset=["Watch Status"],
    )

    st.dataframe(
        styled_df,
        use_container_width=True,
        hide_index=True,
        column_config={
            "Fund": st.column_config.TextColumn(
                "Fund",
                width="large",
            ),
            "Watch Status": st.column_config.TextColumn(
                "Watch Status",
                width="medium",
            ),
        },
    )

    st.markdown("---")

#───Step 15: Single Fund──────────────────────────────────────────────────────────────────

def step15_display_selected_fund():
    import pandas as pd
    import streamlit as st
    import re

    facts = st.session_state.get("fund_factsheets_data", [])
    if not facts:
        st.info("Run Steps 1–14 to populate data before viewing fund details.")
        return

    fund_names = [f["Matched Fund Name"] for f in facts]
    selected_fund = st.selectbox("Select a fund to view details:", fund_names)
    st.session_state.selected_fund = selected_fund

    # =========================================================
    # SELECTED FUND OVERVIEW
    # =========================================================

    st.markdown(f"### {selected_fund}")

    factsheets = st.session_state.get(
        "fund_factsheets_data",
        [],
    )

    factsheet_rec = next(
        (
            row
            for row in factsheets
            if row["Matched Fund Name"] == selected_fund
        ),
        None,
    )

    fund_facts_table = st.session_state.get(
        "step12_fund_facts_table",
        [],
    )

    # Filter out metadata rows if present
    fund_facts_table = [
        row
        for row in fund_facts_table
        if row.get("Fund Name")
        and row.get("Fund Name").lower() != "metadata"
    ]

    # Match the selected fund to its fund-facts record
    facts_rec = next(
        (
            row
            for row in fund_facts_table
            if row.get("Fund Name") == selected_fund
        ),
        None,
    )

    if not facts_rec and factsheet_rec:
        factsheet_ticker = factsheet_rec.get(
            "Matched Ticker"
        )

        facts_rec = next(
            (
                row
                for row in fund_facts_table
                if row.get("Ticker") == factsheet_ticker
            ),
            None,
        )

    if not facts_rec:
        facts_rec = next(
            (
                row
                for row in fund_facts_table
                if selected_fund.lower()
                in row.get("Fund Name", "").lower()
            ),
            None,
        )

    overview_col, statistics_col = st.columns(
        2,
        gap="medium",
    )

    # ---------------------------------------------------------
    # FUND OVERVIEW
    # ---------------------------------------------------------

    with overview_col:
        with st.container(border=True):
            st.markdown("#### Fund Overview")

            if factsheet_rec:
                overview_data = {
                    "Category": factsheet_rec.get(
                        "Category",
                        "—",
                    ),
                    "Benchmark": factsheet_rec.get(
                        "Benchmark",
                        "—",
                    ),
                    "Net Assets": factsheet_rec.get(
                        "Net Assets",
                        "—",
                    ),
                    "Manager": factsheet_rec.get(
                        "Manager Name",
                        "—",
                    ),
                    "Average Market Cap": factsheet_rec.get(
                        "Avg. Market Cap",
                        "—",
                    ),
                }

                for label, value in overview_data.items():
                    row_label, row_value = st.columns(
                        [0.42, 0.58],
                        gap="small",
                    )

                    with row_label:
                        st.markdown(f"**{label}**")

                    with row_value:
                        st.write(value or "—")

            else:
                st.info(
                    "No factsheet information was found "
                    "for this fund."
                )

    # ---------------------------------------------------------
    # FUND STATISTICS
    # ---------------------------------------------------------

    with statistics_col:
        with st.container(border=True):
            st.markdown("#### Fund Statistics")

            if facts_rec:
                statistics_data = {
                    "Manager Tenure": facts_rec.get(
                        "Manager Tenure Yrs.",
                        "—",
                    ),
                    "Expense Ratio": facts_rec.get(
                        "Expense Ratio",
                        "—",
                    ),
                    "Expense Ratio Rank": facts_rec.get(
                        "Expense Ratio Rank",
                        "—",
                    ),
                    "Total Holdings": facts_rec.get(
                        "Total Number of Holdings",
                        "—",
                    ),
                    "Turnover Ratio": facts_rec.get(
                        "Turnover Ratio",
                        "—",
                    ),
                }

                for label, value in statistics_data.items():
                    row_label, row_value = st.columns(
                        [0.55, 0.45],
                        gap="small",
                    )

                    with row_label:
                        st.markdown(f"**{label}**")

                    with row_value:
                        st.write(value or "—")

            else:
                st.info(
                    "No additional fund statistics were found "
                    "for this fund."
                )

    st.markdown("")


    # --- Slide 1 Table: IPS Results ---
    ips_icon_table = st.session_state.get("ips_icon_table")
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected_fund), None)

    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        if not row.empty:
            row_dict = row.iloc[0].to_dict()
            display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
            row_df = pd.DataFrame([{
                "Category": fs_rec.get("Category", "") if fs_rec else "",
                "Time Period": st.session_state.get("report_date", ""),
                "Plan Assets": "$",  # Or replace with actual variable if you store this elsewhere!
                **{display_columns.get(k, k): v for k, v in row_dict.items() if k.startswith("IPS Investment Criteria")},
                "IPS Status": row_dict.get("IPS Watch Status", "")
            }])

            def color_bool(v):
                return "background-color: green" if v == "✔" else ("background-color: red" if v == "✗" else "")
            
            def style_status(v):
                if v == "NW": return "background-color: green; color: white; font-weight: 600;"
                if v == "IW": return "background-color: orange; color: white; font-weight: 600;"
                if v == "FW": return "background-color: red; color: white; font-weight: 600;"
                return ""
            
            styled = row_df.style.applymap(color_bool, subset=[str(i) for i in range(1, 12)]) \
                                 .applymap(style_status, subset=["IPS Status"])
            
            st.dataframe(styled, use_container_width=True)

        else:
            st.warning("No IPS screening result found for selected fund.")
    else:
        st.warning("IPS screening table not found. Run earlier steps first.")

    # --- Slide 2 Table 1: Net Expense Ratio ---
    perf_data = st.session_state.get(
        "fund_performance_data",
        [],
    )

    perf_item = next(
        (
            p
            for p in perf_data
            if p.get("Fund Scorecard Name") == selected_fund
        ),
        {},
    )

    inv_mgr = (
        f"{selected_fund} "
        f"({perf_item.get('Ticker', '')})"
    )

    net_exp = perf_item.get(
        "Net Expense Ratio",
        "",
    )

    if net_exp and not str(net_exp).endswith("%"):
        net_exp = f"{net_exp}%"

    df_slide2 = pd.DataFrame(
        [
            {
                "Fund": inv_mgr,
                "Net Expense Ratio": net_exp,
            }
        ]
    )

  # Save this dataframe for Step 17
    st.session_state["slide2_table1_data"] = df_slide2

    # --- Slide 2 Table 2 ---
    st.markdown("**Returns**")
    date_label = st.session_state.get("report_date", "QTD")
    
    def append_pct(val):
        s = str(val) if val is not None else ""
        return s if s.endswith("%") or s == "" else f"{s}%"
    
    qtd   = append_pct(perf_item.get("QTD", ""))
    one   = append_pct(perf_item.get("1Yr", ""))
    three = append_pct(perf_item.get("3Yr", ""))
    five  = append_pct(perf_item.get("5Yr", ""))
    ten   = append_pct(perf_item.get("10Yr", ""))
    
    bench_qtd   = append_pct(perf_item.get("Bench QTD", ""))
    bench_one   = append_pct(perf_item.get("Bench 1Yr", ""))
    bench_3yr   = append_pct(perf_item.get("Bench 3Yr", ""))
    bench_5yr   = append_pct(perf_item.get("Bench 5Yr", ""))
    bench_ten   = append_pct(perf_item.get("Bench 10Yr", ""))
    
    # Fund row
    row_fund = {
        "Investment Manager": inv_mgr,
        date_label:           qtd,
        "1 Year":             one,
        "3 Year":             three,
        "5 Year":             five,
        "10 Year":            ten
    }
    
    # Benchmark row with corresponding benchmark returns
    bench_name = fs_rec.get("Benchmark", "") if fs_rec else ""
    bench_ticker = fs_rec.get("Matched Ticker", "") if fs_rec else ""
    bench_inv_mgr = f"{bench_name} ({bench_ticker})" if bench_name else "Benchmark"
    
    row_benchmark = {
        "Investment Manager": bench_inv_mgr,
        date_label:           bench_qtd,
        "1 Year":             bench_one,
        "3 Year":             bench_3yr,
        "5 Year":             bench_5yr,
        "10 Year":            bench_ten
    }
    
    df_slide2_2 = pd.DataFrame([row_fund, row_benchmark])
    
    # Save for Step 17 to use
    st.session_state["slide2_table2_data"] = df_slide2_2
    st.dataframe(df_slide2_2, use_container_width=True)


    # --- Slide 2 Table 3 ---
    st.markdown("**Calender Returns**")
    fund_cy = st.session_state.get("step8_returns", [])
    bench_cy = st.session_state.get("benchmark_calendar_year_returns", [])
    if not fund_cy or not bench_cy:
        st.error("❌ No calendar year returns data found. Ensure Step 8 has been run correctly.")
        return
    fund_rec = next((r for r in fund_cy if r.get("Name") == selected_fund), None)
    if not fund_rec:
        st.error(f"❌ Could not find data for selected fund: {selected_fund}")
        return
    benchmark_name = selected_fund
    bench_rec = next((r for r in bench_cy if r.get("Name") == benchmark_name or r.get("Ticker") == fund_rec.get("Ticker")), None)
    if not bench_rec:
        st.error(f"❌ Could not find benchmark data for selected fund: {selected_fund}")
        return
    year_cols = [col for col in fund_rec.keys() if re.match(r"20\d{2}", col)]
    rows = []
    row_fund = {"Investment Manager": f"{selected_fund} ({fund_rec.get('Ticker','')})"}
    for year in year_cols:
        row_fund[year] = fund_rec.get(year, "")
    rows.append(row_fund)
    row_benchmark = {"Investment Manager": f"{bench_rec.get('Name', 'Benchmark')} ({bench_rec.get('Ticker', '')})"}
    for year in year_cols:
        row_benchmark[year] = bench_rec.get(year, "")
    rows.append(row_benchmark)
    df_slide2_3 = pd.DataFrame(rows, columns=["Investment Manager"] + year_cols)

    # Save for Step 17 to use
    st.session_state["slide2_table3_data"] = df_slide2_3
    st.dataframe(df_slide2_3, use_container_width=True)

    # --- Slide 3 Table 1 ---
    st.markdown("**MPT Statistics Summary**")
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    stats3 = next((r for r in mpt3 if r["Fund Name"] == selected_fund), {})
    mpt5 = st.session_state.get("step10_mpt_stats", [])
    stats5 = next((r for r in mpt5 if r["Fund Name"] == selected_fund), {})
    ticker = stats3.get("Ticker", stats5.get("Ticker", ""))
    inv_mgr = f"{selected_fund} ({ticker})"
    row = {
        "Investment Manager":        inv_mgr,
        "3 Year Alpha":              stats3.get("3 Year Alpha", ""),
        "5 Year Alpha":              stats5.get("5 Year Alpha", ""),
        "3 Year Beta":               stats3.get("3 Year Beta", ""),
        "5 Year Beta":               stats5.get("5 Year Beta", ""),
        "3 Year Upside Capture":     stats3.get("3 Year Upside Capture", ""),
        "3 Year Downside Capture":   stats3.get("3 Year Downside Capture", ""),
        "5 Year Upside Capture":     stats5.get("5 Year Upside Capture", ""),
        "5 Year Downside Capture":   stats5.get("5 Year Downside Capture", "")
    }
    df_slide3_1 = pd.DataFrame([row])

    # Save for Step 17 to use
    st.session_state["slide3_table1_data"] = df_slide3_1
    st.dataframe(df_slide3_1, use_container_width=True)

    # --- Slide 3 Table 2 ---
    st.markdown("**Risk-Adjusted Returns / Peer Ranking %**")
    risk_table = st.session_state.get("step13_risk_adjusted_table", [])
    peer_table = st.session_state.get("step14_peer_rank_table", [])
    risk_rec = next((r for r in risk_table if r["Fund Name"] == selected_fund), {})
    peer_rec = next((r for r in peer_table if r["Fund Name"] == selected_fund), {})
    ticker = risk_rec.get("Ticker") or peer_rec.get("Ticker", "")
    inv_mgr = f"{selected_fund} ({ticker})"
    def frac(metric, period):
        r = risk_rec.get(f"{metric} {period}", "")
        p = peer_rec.get(f"{metric} {period}", "")
        return f"{r} / {p}"
    row = {
        "Investment Manager": inv_mgr,
        "3 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "3Yr"),
        "5 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "5Yr"),
        "3 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "3Yr"),
        "5 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "5Yr"),
        "3 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "3Yr"),
        "5 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "5Yr"),
    }
    df_slide3_2 = pd.DataFrame([row])

    # Save for Step 17 to use
    st.session_state["slide3_table2_data"] = df_slide3_2

    st.dataframe(
        df_slide3_2,
        use_container_width=True,
        hide_index=True,
    )
    # --- Slide 4 Table 1: Manager Tenure ---
    blocks = st.session_state.get(
        "fund_blocks",
        [],
    )

    block = next(
        (
            b
            for b in blocks
            if b["Fund Name"] == selected_fund
        ),
        {},
    )

    raw_tenure = next(
        (
            metric["Info"]
            for metric in block.get("Metrics", [])
            if metric["Metric"] == "Manager Tenure"
        ),
        "",
    )

    tenure_match = re.search(
        r"(\d+(?:\.\d+)?)",
        raw_tenure,
    )

    tenure = (
        f"{tenure_match.group(1)} years"
        if tenure_match
        else raw_tenure
    )

    perf_data = st.session_state.get(
        "fund_performance_data",
        [],
    )

    perf_item = next(
        (
            p
            for p in perf_data
            if p.get("Fund Scorecard Name") == selected_fund
        ),
        {},
    )

    inv_mgr = (
        f"{selected_fund} "
        f"({perf_item.get('Ticker', '')})"
    )

    df_slide4 = pd.DataFrame(
        [
            {
                "Investment Manager": inv_mgr,
                "Manager Tenure": tenure,
            }
        ]
    )

    # Keep this unchanged for PowerPoint export.
    st.session_state["slide4"] = df_slide4

    # Net Expense Ratio and Manager Tenure are already shown in
    # the summary cards above. Their DataFrames remain saved in
    # session state for PowerPoint export, but are not displayed here.

    # --- Slide 4 Table 2: Assets ---
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected_fund), None)
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p["Fund Scorecard Name"] == selected_fund), None)
    inv_mgr    = f"{selected_fund} ({perf_item.get('Ticker','') if perf_item else ''})"
    assets     = fs_rec.get("Net Assets", "") if fs_rec else ""
    avg_cap    = fs_rec.get("Avg. Market Cap", "") if fs_rec else ""
    df_slide4_2 = pd.DataFrame([{
        "Investment Manager":             inv_mgr,
        "Assets Under Management":        assets,
        "Average Market Capitalization":  avg_cap
    }])

    # Save for Step 17 to use.
    # Assets are already shown in the Fund Overview card, so this
    # table is kept for PowerPoint export but not displayed here.
    st.session_state["slide4_table2_data"] = df_slide4_2

#───Bullet Points──────────────────────────────────────────────────────────────────

def step16_bullet_points():
    import streamlit as st

    selected_fund = st.session_state.get("selected_fund")
    if not selected_fund:
        st.error("❌ No fund selected. Please select a fund in Step 15.")
        return

    perf_data = st.session_state.get("fund_performance_data", [])
    item = next((x for x in perf_data if x["Fund Scorecard Name"] == selected_fund), None)
    if not item:
        st.error(f"❌ Performance data for '{selected_fund}' not found.")
        return

    bullets = []

    # Bullet 1: Performance vs. Benchmark, using template
    template = st.session_state.get("bullet_point_templates", [""])[0]
    b1 = template
    for fld, val in item.items():
        b1 = b1.replace(f"[{fld}]", str(val))
    bullets.append(b1)
    st.markdown(f"- {b1}")

    # Get IPS status from icon table (guaranteed to match Slide 1 Table)
    ips_icon_table = st.session_state.get("ips_icon_table")
    ips_status = None
    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        ips_status = row.iloc[0]["IPS Watch Status"] if not row.empty else None

    # Bullet 2: Watch status and return comparison
    if ips_status == "NW":
        b2 = "- This fund is **not on watch**."
        bullets.append(b2)
        st.markdown(b2)
    else:
        status_label = (
            "Formal Watch" if ips_status == "FW" else
            "Informal Watch" if ips_status == "IW" else
            ips_status or "on watch"
        )

        three   = float(item.get("3Yr") or 0)
        bench3  = float(item.get("Bench 3Yr") or 0)
        five    = float(item.get("5Yr") or 0)
        bench5  = float(item.get("Bench 5Yr") or 0)
        bps3 = round((three - bench3) * 100, 1)
        bps5 = round((five  - bench5) * 100, 1)

        # Peer rank logic (safe handling)
        peer = st.session_state.get("step14_peer_rank_table", [])
        raw3 = next((r.get("Sharpe Ratio 3Yr") or r.get("Sharpe Ratio Rank 3Yr") for r in peer
                     if r.get("Fund Name") == selected_fund), None)
        raw5 = next((r.get("Sharpe Ratio 5Yr") or r.get("Sharpe Ratio Rank 5Yr") for r in peer
                     if r.get("Fund Name") == selected_fund), None)
        try:
            pos3 = "top" if raw3 and int(raw3) <= 50 else "bottom"
        except:
            pos3 = "bottom"
        try:
            pos5 = "top" if raw5 and int(raw5) <= 50 else "bottom"
        except:
            pos5 = "bottom"

        b2 = (
            f"- The fund is now on **{status_label}**. Its 3‑year return trails the benchmark by "
            f"{bps3} bps ({three:.2f}% vs. {bench3:.2f}%) and its 5‑year return trails by "
            f"{bps5} bps ({five:.2f}% vs. {bench5:.2f}%). Its 3‑Yr Sharpe ranks in the {pos3} half of peers "
            f"and its 5‑Yr Sharpe ranks in the {pos5} half."
        )
        bullets.append(b2)
        st.markdown(b2)

    # Bullet 3: Action for Formal Watch only
    if ips_status == "FW":
        b3 = "- **Action:** Consider replacing this fund."
        bullets.append(b3)
        st.markdown(b3)

    # Save bullets for Step 17
    st.session_state["bullet_points"] = bullets


#───Build Powerpoint──────────────────────────────────────────────────────────────────
def step17_export_to_ppt():
    import streamlit as st
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from io import BytesIO
    import pandas as pd

    selected = st.session_state.get("selected_fund")
    if not selected:
        st.error("❌ No fund selected. Please select a fund in Step 15.")
        return

    template_path = "assets/writeup_templates.pptx"  # Adjust path if needed
    try:
        prs = Presentation(template_path)
    except Exception as e:
        st.error(f"Could not load PowerPoint template: {e}")
        return

    # Helper to get table header texts
    def get_table_header(table):
        return tuple(cell.text.strip() for cell in table.rows[0].cells)
   
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.util import Pt
    
    def fill_table_with_styles(table, df_table, bold_row_idx=None, first_col_white=True):
        """
        Fill a PPT table with df_table values, styling:
          - First column font: white or black
          - ✔ cells green text, ✗ cells red text
          - IPS Status cell background per status (NW/IW/FW)
          - Bold the specified data row
        """
        # Map IPS statuses to background colors
        status_bg = {
            "NW": RGBColor(0xC5, 0xE6, 0x9A),  # green
            "IW": RGBColor(0xFF, 0x95, 0x53),  # orange
            "FW": RGBColor(0xFF, 0x5D, 0x58),  # red
        }
    
        n_rows = min(len(df_table), len(table.rows) - 1)
        n_cols = min(len(df_table.columns), len(table.columns))
    
        # Find the column index for "IPS Status"
        try:
            status_col = list(df_table.columns).index("IPS Status")
        except ValueError:
            status_col = None
    
        for i in range(n_rows):
            for j in range(n_cols):
                val = df_table.iloc[i, j]
                cell = table.cell(i + 1, j)  # data rows start at row 1
                cell.text = str(val) if val is not None else ""
                cell.vertical_alignment = MSO_VERTICAL_ANCHOR.MIDDLE
    
                # If this is the IPS Status column on the (first) data row, color its background
                if status_col is not None and j == status_col and i == 0:
                    cell.fill.solid()
                    bg_color = status_bg.get(str(val), None)
                    if bg_color:
                        cell.fill.fore_color.rgb = bg_color
    
                # Now style the text runs
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.name = "Cambria"
                        run.font.size = Pt(11)
    
                        # first column font color
                        if j == 0:
                            run.font.color.rgb = RGBColor(255,255,255) if first_col_white else RGBColor(0,0,0)
                        else:
                            # green ✔, red ✗, else black
                            if run.text == "✔":
                                run.font.color.rgb = RGBColor(0,128,0)
                            elif run.text == "✗":
                                run.font.color.rgb = RGBColor(255,0,0)
                            else:
                                run.font.color.rgb = RGBColor(0,0,0)
    
                        # bold entire row if requested
                        run.font.bold = (bold_row_idx is not None and i == bold_row_idx)


    # Replace placeholder text in slide, preserving formatting as best as possible
    def fill_text_placeholder_preserving_format(slide, placeholder_text, replacement_text):
        replaced = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs)
                if placeholder_text in full_text:
                    new_text = full_text.replace(placeholder_text, replacement_text)
                    runs = paragraph.runs
                    if len(runs) == 1:
                        runs[0].text = new_text
                    else:
                        for run in runs:
                            run.text = ""
                        avg_len = len(new_text) // len(runs)
                        idx = 0
                        for run in runs[:-1]:
                            run.text = new_text[idx:idx+avg_len]
                            idx += avg_len
                        runs[-1].text = new_text[idx:]
                    replaced = True
        return replaced

    # Fill bullet points placeholder with a list of bullet strings
    def fill_bullet_points(slide, placeholder="[Bullet Point 1]", bullets=None):
        import streamlit as st
        if bullets is None:
            bullets = st.session_state.get("bullet_points", None)
        if not bullets:
            bullets = ["Performance exceeded benchmark.", "No watch status.", "No action required."]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                if placeholder in p.text:
                    shape.text_frame.clear()
                    for b in bullets:
                        p = shape.text_frame.add_paragraph()
                        # Remove markdown asterisks if any
                        clean_text = b.replace("**", "")
                        p.text = clean_text
                        p.level = 0
                        p.font.name = "Cambria"
                        p.font.size = Pt(11)
                        p.font.color.rgb = RGBColor(0, 0, 0)
                        # Bold the entire paragraph if you want
                        p.font.bold = True  # <-- makes the whole bullet bold
                    return True
        return False


    # Get actual category for placeholder replacement
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected), {})
    category = fs_rec.get("Category", "N/A")

    # Prepare Slide data from session_state
    df_slide1 = None
    ips_icon_table = st.session_state.get("ips_icon_table")
    if ips_icon_table is not None and not ips_icon_table.empty:
        filtered = ips_icon_table[ips_icon_table["Fund Name"] == selected]
        if not filtered.empty:
            row = filtered.iloc[0]
            display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
            table_data = {
                **{display_columns.get(k, k): v for k, v in row.items() if k.startswith("IPS Investment Criteria")},
                "IPS Status": row.get("IPS Watch Status", "")
            }
            table_data["Category"] = category
            table_data["Time Period"] = st.session_state.get("report_date", "")
            table_data["Plan Assets"] = "$"
            headers = ["Category", "Time Period", "Plan Assets"] + [str(i+1) for i in range(11)] + ["IPS Status"]
            df_slide1 = pd.DataFrame([table_data], columns=headers)

    df_slide2_table1 = st.session_state.get("slide2_table1_data")
    df_slide2_table2 = st.session_state.get("slide2_table2_data")
    df_slide2_table3 = st.session_state.get("slide2_table3_data")

    df_slide3_table1 = st.session_state.get("slide3_table1_data")
    df_slide3_table2 = st.session_state.get("slide3_table2_data")

    df_slide4_table1 = st.session_state.get("slide4")
    df_slide4_table2 = st.session_state.get("slide4_table2_data")

    # --- Fill Slide 1 ---
    slide1 = prs.slides[0]
    if not fill_text_placeholder_preserving_format(slide1, "[Fund Name]", selected):
        st.warning("Could not find the [Fund Name] placeholder on Slide 1.")

    if df_slide1 is not None:
        filled = False
        for shape in slide1.shapes:
            if shape.has_table:
                table = shape.table
                if len(table.columns) == len(df_slide1.columns):
                    # Use first_col_white=False here to set first column font black for Slide 1 table
                    fill_table_with_styles(table, df_slide1, first_col_white=False)
                    filled = True
                    break
        if not filled:
            st.warning("Could not find matching table on Slide 1 to fill.")
    else:
        st.warning("Slide 1 IPS data not found in session state.")

    # Fill bullet points on Slide 1
    bullets = st.session_state.get("bullet_points", None)
    if not fill_bullet_points(slide1, "[Bullet Point 1]", bullets):
        st.warning("Could not find bullet points placeholder on Slide 1.")

    # --- Fill Slide 2 ---
    slide2 = prs.slides[1]
    if not fill_text_placeholder_preserving_format(slide2, "[Category]", category):
        st.warning("Could not find [Category] placeholder on Slide 2.")

    if df_slide2_table1 is None:
        st.warning("Slide 2 Table 1 data not found.")
    else:
        filled = False
        for shape in slide2.shapes:
            if shape.has_table:
                table = shape.table
                if len(table.columns) == len(df_slide2_table1.columns):
                    fill_table_with_styles(table, df_slide2_table1)  # default: first_col_white=True
                    filled = True
                    break
        if not filled:
            st.warning("Could not find matching table for Slide 2 Table 1.")

    if df_slide2_table2 is None:
        st.warning("Slide 2 Table 2 data not found.")
    else:
        # Replace 2nd header cell text with quarter/year label
        quarter_label = st.session_state.get("report_date", "QTD")
        filled = False
        for shape in slide2.shapes:
            if shape.has_table:
                table = shape.table
                if len(table.columns) == len(df_slide2_table2.columns):
                    # Update header row font/color + replace 2nd header cell
                    table.cell(0, 1).text = quarter_label
                    for c in range(len(table.columns)):
                        cell = table.cell(0, c)
                        for para in cell.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            for run in para.runs:
                                run.font.name = "Cambria"
                                run.font.size = Pt(11)
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True
                    fill_table_with_styles(table, df_slide2_table2, bold_row_idx=1)  # bold benchmark row, first_col_white=True
                    filled = True
                    break
        if not filled:
            st.warning("Could not find matching table for Slide 2 Table 2.")

    if df_slide2_table3 is None:
        st.warning("Slide 2 Table 3 data not found.")
    else:
        # Update headers dynamically for years, fill benchmark row bold
        filled = False
        for shape in slide2.shapes:
            if shape.has_table:
                table = shape.table
                if len(table.columns) == len(df_slide2_table3.columns):
                    # Replace header cells with df column names & style font/color
                    for c, col in enumerate(df_slide2_table3.columns):
                        cell = table.cell(0, c)
                        cell.text = str(col)
                        for para in cell.text_frame.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            for run in para.runs:
                                run.font.name = "Cambria"
                                run.font.size = Pt(11)
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True
                    fill_table_with_styles(table, df_slide2_table3, bold_row_idx=1)  # bold benchmark row, first_col_white=True
                    filled = True
                    break
        if not filled:
            st.warning("Could not find matching table for Slide 2 Table 3.")

    # --- Fill Slide 3 ---
    slide3 = prs.slides[2]
    if not fill_text_placeholder_preserving_format(slide3, "[Category]", category):
        st.warning("Could not find [Category] placeholder on Slide 3.")

    if df_slide3_table1 is None or df_slide3_table2 is None:
        st.warning("Slide 3 table data not found.")
    else:
        tables = [shape.table for shape in slide3.shapes if shape.has_table]
        if len(tables) < 2:
            st.warning("Expected two tables on Slide 3, found fewer.")
        else:
            if len(df_slide3_table1.columns) == len(tables[0].columns):
                fill_table_with_styles(tables[0], df_slide3_table1)  # first_col_white=True default
            else:
                st.warning("Slide 3 Table 1 headers do not match.")
            if len(df_slide3_table2.columns) == len(tables[1].columns):
                fill_table_with_styles(tables[1], df_slide3_table2)  # first_col_white=True default
            else:
                st.warning("Slide 3 Table 2 headers do not match.")

    # --- Fill Slide 4 ---
    slide4 = prs.slides[3]
    qualitative_placeholder = f"[Category]– Qualitative Factors"
    qualitative_replacement = f"{category} - Qualitative Factors"
    if not fill_text_placeholder_preserving_format(slide4, qualitative_placeholder, qualitative_replacement):
        st.warning(f"Could not find placeholder '{qualitative_placeholder}' on Slide 4.")

    if df_slide4_table1 is not None:
        filled = False
        for shape in slide4.shapes:
            if shape.has_table:
                table = shape.table
                if len(table.columns) == len(df_slide4_table1.columns):
                    fill_table_with_styles(table, df_slide4_table1)  # first_col_white=True default
                    filled = True
                    break
        if not filled:
            st.warning("Could not find matching table for Slide 4 Table 1.")

    if df_slide4_table2 is not None:
        filled = False
        for shape in slide4.shapes:
            if shape.has_table:
                table = shape.table
                if len(table.columns) == len(df_slide4_table2.columns):
                    fill_table_with_styles(table, df_slide4_table2)  # first_col_white=True default
                    filled = True
                    break
        if not filled:
            st.warning("Could not find matching table for Slide 4 Table 2.")

    # --- Save and offer download ---
    output = BytesIO()
    prs.save(output)
    st.success("Powerpoint Generated")
    st.download_button(
        label="Download Writeup PowerPoint",
        data=output.getvalue(),
        file_name=f"{selected} Writeup.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ─── Main App ────────────────────────────────────────────────────────────────

def run():
    import re

    st.set_page_config(
        page_title="Writeup",
        layout="wide",
    )

    st.markdown(
        """
        <style>
            /* =========================================================
               PAGE
               ========================================================= */

            .block-container {
                max-width: 1450px;
                padding-top: 2rem;
                padding-bottom: 3rem;
            }

            h1 {
                color: #16243A;
                letter-spacing: -0.03em;
                margin-bottom: 0.3rem;
            }

            h2,
            h3 {
                color: #24364B;
                letter-spacing: -0.02em;
            }

            .app-subtitle {
                color: #667085;
                font-size: 1rem;
                line-height: 1.55;
                margin-bottom: 1.35rem;
                max-width: 900px;
            }

            .section-header {
                margin-top: 0.4rem;
                margin-bottom: 1rem;
            }

            .section-label {
                color: #667085;
                font-size: 0.74rem;
                font-weight: 700;
                letter-spacing: 0.09em;
                text-transform: uppercase;
                margin-bottom: 0.25rem;
            }

            .section-title {
                color: #24364B;
                font-size: 1.1rem;
                font-weight: 700;
                letter-spacing: -0.015em;
                margin-bottom: 0.25rem;
            }

            .section-description {
                color: #667085;
                font-size: 0.87rem;
                line-height: 1.5;
                max-width: 900px;
            }

            .page-divider {
                height: 1px;
                background: #E4E7EC;
                margin-bottom: 1.35rem;
            }

            /* =========================================================
               FILE UPLOADER
               ========================================================= */

            div[data-testid="stFileUploader"] {
                border: 1px solid #D0D5DD;
                border-radius: 12px;
                background: #FFFFFF;
                padding: 0.4rem;
                margin-bottom: 0.8rem;
            }

            div[data-testid="stFileUploader"] section {
                background: #FAFBFC;
                border: 1px dashed #C8D0DC;
                border-radius: 9px;
            }

            div[data-testid="stFileUploader"]:hover {
                border-color: #AEB9C6;
            }

            /* =========================================================
               PROGRESS
               ========================================================= */

            div[data-testid="stProgress"] {
                margin-top: 0.4rem;
                margin-bottom: 0.2rem;
            }

            div[data-testid="stProgress"] > div > div {
                border-radius: 999px;
            }

            /* =========================================================
               MAIN TABS
               ========================================================= */

            div[data-testid="stTabs"]
            > div[data-baseweb="tab-list"] {
                gap: 1.5rem;
                border-bottom: 1px solid #D9DEE7;
                margin-top: 0.5rem;
                margin-bottom: 1.25rem;
            }

            div[data-testid="stTabs"]
            button[data-baseweb="tab"] {
                background: transparent;
                border: none;
                border-radius: 0;
                color: #667085;
                font-size: 0.88rem;
                font-weight: 600;
                padding: 0.7rem 0.05rem 0.75rem 0.05rem;
            }

            div[data-testid="stTabs"]
            button[data-baseweb="tab"]:hover {
                color: #24364B;
                background: transparent;
            }

            div[data-testid="stTabs"]
            button[data-baseweb="tab"][aria-selected="true"] {
                color: #16243A;
                font-weight: 700;
            }

            div[data-testid="stTabs"]
            div[data-baseweb="tab-highlight"] {
                background-color: #162F52;
                height: 2px;
            }

            div[data-testid="stTabs"]
            div[data-baseweb="tab-border"] {
                background-color: transparent;
            }

            div[data-testid="stTabPanel"] {
                padding-top: 0.2rem;
            }

            /* =========================================================
               EXPANDERS
               ========================================================= */

            div[data-testid="stExpander"] {
                background: #FFFFFF;
                border: 1px solid #E4E7EC;
                border-radius: 10px;
                overflow: hidden;
                box-shadow: none;
                margin-bottom: 0.65rem;
            }

            div[data-testid="stExpander"] summary {
                min-height: 47px;
            }

            div[data-testid="stExpander"] summary:hover {
                background: #F8FAFC;
            }

            div[data-testid="stExpander"] summary p {
                color: #344054;
                font-size: 0.9rem;
                font-weight: 650;
            }

            div[data-testid="stExpander"] details[open] > summary {
                background: #FAFBFC;
                border-bottom: 1px solid #EAECF0;
            }

            /* =========================================================
               METRICS AND TABLES
               ========================================================= */

            div[data-testid="stMetric"] {
                background: #FFFFFF;
                border: 1px solid #E4E7EC;
                border-radius: 10px;
                padding: 0.85rem 1rem;
                box-shadow: none;
            }

            div[data-testid="stMetricLabel"] {
                color: #667085;
                font-size: 0.8rem;
                font-weight: 600;
            }

            div[data-testid="stMetricValue"] {
                color: #16243A;
                font-weight: 700;
            }

            div[data-testid="stDataFrame"],
            div[data-testid="stTable"] {
                border: 1px solid #E4E7EC;
                border-radius: 10px;
                overflow: hidden;
                box-shadow: none;
            }

            /* =========================================================
               BUTTONS
               ========================================================= */

            .stButton > button,
            .stDownloadButton > button {
                border-radius: 8px;
                border: 1px solid #CDD5DF;
                background: #FFFFFF;
                color: #344054;
                font-weight: 620;
                min-height: 39px;
                padding-left: 1rem;
                padding-right: 1rem;
                box-shadow: none;
            }

            .stButton > button:hover,
            .stDownloadButton > button:hover {
                background: #F8FAFC;
                border-color: #AEB9C6;
                color: #24364B;
            }

            .stButton > button:active,
            .stDownloadButton > button:active {
                background: #F2F4F7;
            }

            /* =========================================================
               INPUTS
               ========================================================= */

            div[data-baseweb="select"] > div,
            div[data-baseweb="input"] > div,
            div[data-baseweb="textarea"] > div,
            textarea {
                border-radius: 8px !important;
            }

            /* =========================================================
               ALERTS
               ========================================================= */

            div[data-testid="stAlert"] {
                border-radius: 9px;
                border-width: 1px;
            }
        </style>
        """,
        unsafe_allow_html=True,
    )

    # =========================================================
    # PAGE HEADER
    # =========================================================

    st.title("Writeup")

    st.markdown(
        """
        <div class="app-subtitle">
            Review MPI fund data, evaluate investment criteria, prepare fund
            commentary, and generate the completed PowerPoint presentation.
        </div>
        <div class="page-divider"></div>
        """,
        unsafe_allow_html=True,
    )

    # =========================================================
    # REPORT UPLOAD
    # =========================================================

    st.markdown(
        """
        <div class="section-header">
            <div class="section-label">Source Report</div>
            <div class="section-title">Upload MPI PDF</div>
            <div class="section-description">
                Select the report that will be used throughout the fund review,
                writeup, and presentation workflow.
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    uploaded = st.file_uploader(
        "Upload MPI PDF",
        type=["pdf"],
        help="Upload a text-based MPI report PDF.",
        key="writeup_generator_pdf_uploader",
        label_visibility="collapsed",
    )

    if not uploaded:
        st.caption("Upload an MPI PDF to begin processing.")
        return

    # =========================================================
    # PROGRESS DISPLAY
    # =========================================================

    progress_bar = st.progress(0)
    progress_message = st.empty()

    def update_progress(percent, message):
        progress_message.caption(message)
        progress_bar.progress(percent)

    update_progress(
        2,
        "Opening report...",
    )

    with pdfplumber.open(uploaded) as pdf:
        # =========================================================
        # STEP 1: REPORT INFORMATION
        # =========================================================

        update_progress(
            8,
            "Reading report information...",
        )

        first = pdf.pages[0].extract_text() or ""
        process_page1(first)

        # =========================================================
        # STEP 2: TABLE OF CONTENTS
        # =========================================================

        update_progress(
            15,
            "Reading report structure...",
        )

        toc_text = "".join(
            (pdf.pages[i].extract_text() or "")
            for i in range(min(3, len(pdf.pages)))
        )

        process_toc(toc_text)

        # =========================================================
        # MAIN PAGE TABS
        # =========================================================

        (
            overview_tab,
            analysis_tab,
            writeup_tab,
            export_tab,
        ) = st.tabs(
            [
                "Report Overview",
                "Fund Analysis",
                "Writeup",
                "Export",
            ]
        )

        # =========================================================
        # REPORT OVERVIEW TAB
        # =========================================================

        with overview_tab:
            st.markdown(
                """
                <div class="section-header">
                    <div class="section-label">Report Overview</div>
                    <div class="section-title">Report Information</div>
                    <div class="section-description">
                        Review the report details identified during processing.
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )

            show_report_summary()
           
            # -----------------------------------------------------
            # IPS FAIL TABLE
            # -----------------------------------------------------

            step14_5_ips_fail_table()
          
            with st.expander(
                "Table of Contents",
                expanded=False,
            ):
                st.write(
                    "The report structure has been identified and saved "
                    "for the fund analysis process."
                )

                toc_details = {
                    "Fund Scorecard Page": st.session_state.get(
                        "scorecard_page"
                    ),
                    "Performance Page": st.session_state.get(
                        "performance_page"
                    ),
                    "Calendar-Year Returns Page": st.session_state.get(
                        "calendar_year_page"
                    ),
                    "3-Year MPT Page": st.session_state.get(
                        "r3yr_page"
                    ),
                    "5-Year MPT Page": st.session_state.get(
                        "r5yr_page"
                    ),
                    "Fund Factsheets Page": st.session_state.get(
                        "factsheets_page"
                    ),
                }

                for label, page_number in toc_details.items():
                    displayed_page = (
                        page_number
                        if page_number is not None
                        else "Not included"
                    )

                    st.write(
                        f"**{label}:** {displayed_page}"
                    )

        # =========================================================
        # FUND ANALYSIS TAB
        # =========================================================

        with analysis_tab:
            st.markdown(
                """
                <div class="section-header">
                    <div class="section-label">Fund Analysis</div>
                    <div class="section-title">Extracted Fund Details</div>
                    <div class="section-description">
                        Review the screening, factsheet, returns, and risk
                        information extracted from the report.
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )

            # -----------------------------------------------------
            # IPS INVESTMENT SCREENING
            # -----------------------------------------------------

            update_progress(
                25,
                "Processing IPS investment screening...",
            )

            with st.expander(
                "IPS Investment Screening",
                expanded=True,
            ):
                sp = st.session_state.get(
                    "scorecard_page"
                )

                tot = st.session_state.get(
                    "total_options"
                )

                pp = st.session_state.get(
                    "performance_page"
                )

                factsheets_page = st.session_state.get(
                    "factsheets_page"
                )

                if sp and tot is not None and pp:
                    step3_5_6_scorecard_and_ips(
                        pdf,
                        sp,
                        pp,
                        factsheets_page,
                        tot,
                    )
                else:
                    st.info(
                        "IPS screening information was not found "
                        "in this report."
                    )

            # -----------------------------------------------------
            # FUND FACTSHEETS
            # -----------------------------------------------------

            update_progress(
                42,
                "Processing fund factsheets...",
            )

            with st.expander(
                "Fund Factsheets",
                expanded=False,
            ):
                names = [
                    block["Fund Name"]
                    for block in st.session_state.get(
                        "fund_blocks",
                        [],
                    )
                ]

                if names:
                    step6_process_factsheets(
                        pdf,
                        names,
                    )

                    factsheet_data = st.session_state.get(
                        "fund_factsheets_data",
                        [],
                    )

                    if factsheet_data:
                        st.success(
                            f"Processed {len(factsheet_data)} "
                            "fund factsheets."
                        )
                    else:
                        st.info(
                            "No fund factsheets were found "
                            "in this report."
                        )
                else:
                    st.info(
                        "No funds were available for factsheet processing."
                    )

            # -----------------------------------------------------
            # FUND FACTS
            # -----------------------------------------------------

            update_progress(
                52,
                "Extracting fund facts...",
            )

            with st.expander(
                "Fund Facts",
                expanded=False,
            ):
                step12_process_fund_facts(pdf)

                fund_facts = st.session_state.get(
                    "step12_fund_facts_table",
                    [],
                )

                if fund_facts:
                    st.success(
                        f"Fund facts were extracted for "
                        f"{len(fund_facts)} funds."
                    )
                else:
                    st.info(
                        "No additional fund facts were found."
                    )

            # -----------------------------------------------------
            # RETURNS
            # -----------------------------------------------------

            update_progress(
                62,
                "Processing fund returns...",
            )

            with st.expander(
                "Returns",
                expanded=False,
            ):
                step7_extract_returns(pdf)
                step8_calendar_returns(pdf)

            # -----------------------------------------------------
            # MPT STATISTICS
            # -----------------------------------------------------

            update_progress(
                72,
                "Checking for MPT statistics...",
            )

            with st.expander(
                "MPT Statistics Summary",
                expanded=False,
            ):
                has_mpt_3yr = (
                    st.session_state.get("r3yr_page")
                    is not None
                )

                has_mpt_5yr = (
                    st.session_state.get("r5yr_page")
                    is not None
                )

                if has_mpt_3yr:
                    step9_risk_analysis_3yr(pdf)

                if has_mpt_5yr:
                    step10_risk_analysis_5yr(pdf)

                mpt3 = st.session_state.get(
                    "step9_mpt_stats",
                    [],
                )

                mpt5 = st.session_state.get(
                    "step10_mpt_stats",
                    [],
                )

                if mpt3 and mpt5:
                    step11_create_summary()

                elif mpt3 and not mpt5:
                    st.info(
                        "Three-year MPT statistics were included, "
                        "but five-year statistics were not available."
                    )

                    st.dataframe(
                        pd.DataFrame(mpt3),
                        use_container_width=True,
                        hide_index=True,
                    )

                elif mpt5 and not mpt3:
                    st.info(
                        "Five-year MPT statistics were included, "
                        "but three-year statistics were not available."
                    )

                    st.dataframe(
                        pd.DataFrame(mpt5),
                        use_container_width=True,
                        hide_index=True,
                    )

                else:
                    st.info(
                        "No MPT statistics were included in this report."
                    )

            # -----------------------------------------------------
            # RISK-ADJUSTED RETURNS
            # -----------------------------------------------------

            update_progress(
                82,
                "Processing risk-adjusted returns...",
            )

            with st.expander(
                "Risk-Adjusted Returns",
                expanded=False,
            ):
                step13_process_risk_adjusted_returns(pdf)
                step14_extract_peer_risk_adjusted_return_rank(
                    pdf
                )

        # =========================================================
        # DATA PREPARATION FOR WRITEUP
        # =========================================================

        update_progress(
            88,
            "Preparing writeup data...",
        )

        report_date = st.session_state.get(
            "report_date",
            "",
        )

        match = re.match(
            r"(\d)(?:st|nd|rd|th)\s+QTR,\s*(\d{4})",
            report_date or "",
        )

        quarter = match.group(1) if match else ""
        year = match.group(2) if match else ""

        for itm in st.session_state.get(
            "fund_performance_data",
            [],
        ):
            qtd = float(
                itm.get("QTD") or 0
            )

            bench_qtd = float(
                itm.get("Bench QTD") or 0
            )

            itm["Perf Direction"] = (
                "overperformed"
                if qtd >= bench_qtd
                else "underperformed"
            )

            itm["Quarter"] = quarter
            itm["Year"] = year

            diff_bps = round(
                (qtd - bench_qtd) * 100,
                1,
            )

            itm["QTD_bps_diff"] = str(
                diff_bps
            )

            fund_pct = f"{qtd:.2f}%"
            bench_pct = f"{bench_qtd:.2f}%"

            itm["QTD_pct_diff"] = (
                f"{(qtd - bench_qtd):.2f}%"
            )

            itm["QTD_vs"] = (
                f"{fund_pct} vs. {bench_pct}"
            )

        if "bullet_point_templates" not in st.session_state:
            st.session_state[
                "bullet_point_templates"
            ] = [
                "[Fund Scorecard Name] [Perf Direction] its benchmark in Q[Quarter], "
                "[Year] by [QTD_bps_diff] bps ([QTD_vs])."
            ]

        # =========================================================
        # WRITEUP TAB
        # =========================================================
        with writeup_tab:
        
            st.markdown(
                """
                <div class="section-header">
                    <div class="section-label">WRITEUP</div>
                    <div class="section-title">
                        Review and Prepare Commentary
                    </div>
                    <div class="section-description">
                        Review the selected fund and prepare the commentary
                        for the presentation.
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )
        
            step15_display_selected_fund()
        
            st.markdown("---")
        
            with st.expander(
                "Bullet Points",
                expanded=False,
            ):
                step16_bullet_points()

        # =========================================================
        # EXPORT TAB
        # =========================================================

        update_progress(
            96,
            "Preparing presentation export...",
        )

        with export_tab:
            st.markdown(
                """
                <div class="section-header">
                    <div class="section-label">Final Output</div>
                    <div class="section-title">PowerPoint Presentation</div>
                    <div class="section-description">
                        Review the final selections and generate the completed
                        fund review presentation.
                    </div>
                </div>
                """,
                unsafe_allow_html=True,
            )

            step17_export_to_ppt()

        # =========================================================
        # COMPLETE
        # =========================================================

        progress_bar.progress(100)
        progress_message.success(
            "Report processing complete."
        )


if __name__ == "__main__":
    run()
