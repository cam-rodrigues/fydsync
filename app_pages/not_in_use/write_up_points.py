import re
import streamlit as st
import pdfplumber
from calendar import month_name
import pandas as pd
from rapidfuzz import fuzz
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import yfinance as yf

def extract_performance_table(pdf, performance_page, fund_names, end_page=None):
    import re
    from rapidfuzz import fuzz

    # Decide where to stop in the PDF
    end = end_page if end_page is not None else len(pdf.pages) + 1

    # 1. Get all lines from the section
    lines = []
    for pnum in range(performance_page - 1, end - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]

    # 2. Prepare regex
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")

    # 3. For each fund, try to pull numbers
    perf_data = []
    for name in fund_names:
        item = {"Fund Scorecard Name": name}
        tk = ""  # You’ll fill in ticker later from st.session_state["tickers"] or similar
        # a) Exact-ticker match: you’ll want to match this if you have tickers already
        idx = next(
            (i for i, ln in enumerate(lines)
             if name in ln),
            None
        )
        # b) Fuzzy-name fallback if not found
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower()))
                      for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                continue  # Can't find a match

        # c) Pull fund numbers from line above (and two above if needed)
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        clean += [None] * (8 - len(clean))  # pad

        # d) Map to columns
        item["QTD"] = clean[0]
        item["1Yr"] = clean[2]
        item["3Yr"] = clean[3]
        item["5Yr"] = clean[4]
        item["10Yr"] = clean[5]
        item["Net Expense Ratio"] = clean[-2]

        # e) Pull benchmark QTD, 3Yr, 5Yr from next lines
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 1 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]
        item["Bench QTD"] = bench_clean[0] if bench_clean else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None

        perf_data.append(item)
    return perf_data


#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
# === Utility: Extract & Label Report Date ===
def extract_report_date(text):
    # find the first quarter‐end or any mm/dd/yyyy
    dates = re.findall(r'(\d{1,2})/(\d{1,2})/(20\d{2})', text or "")
    for month, day, year in dates:
        m, d = int(month), int(day)
        # quarter‐end mapping
        if (m, d) in [(3,31), (6,30), (9,30), (12,31)]:
            q = { (3,31): "1st", (6,30): "2nd", (9,30): "3rd", (12,31): "4th" }[(m,d)]
            return f"{q} QTR, {year}"
        # fallback: human‐readable
        return f"As of {month_name[m]} {d}, {year}"
    return None
#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# === Step 1 & 1.5: Page 1 Extraction ===
def process_page1(text):
    report_date = extract_report_date(text)
    if report_date:
        st.session_state['report_date'] = report_date
        st.success(f"Report Date: {report_date}")
    else:
        st.error("Could not detect report date on page 1.")

    m = re.search(r"Total Options:\s*(\d+)", text or "")
    st.session_state['total_options'] = int(m.group(1)) if m else None

    m = re.search(r"Prepared For:\s*\n(.*)", text or "")
    st.session_state['prepared_for'] = m.group(1).strip() if m else None

    m = re.search(r"Prepared By:\s*(.*)", text or "")
    pb = m.group(1).strip() if m else ""
    if not pb or "mpi stylus" in pb.lower():
        pb = "Procyon Partners, LLC"
    st.session_state['prepared_by'] = pb

    st.subheader("Page 1 Metadata")
    st.write(f"- Total Options: {st.session_state['total_options']}")
    st.write(f"- Prepared For: {st.session_state['prepared_for']}")
    st.write(f"- Prepared By: {pb}")


# === Step 2: Table of Contents Extraction ===
def process_toc(text):
    perf = re.search(r"Fund Performance[^\d]*(\d{1,3})", text or "")
    sc   = re.search(r"Fund Scorecard\s+(\d{1,3})", text or "")
    fs   = re.search(r"Fund Factsheets\s+(\d{1,3})", text or "")
    cy   = re.search(r"Fund Performance: Calendar Year\s+(\d{1,3})", text or "")
    r3yr = re.search(r"Risk Analysis: MPT Statistics \(3Yr\)\s+(\d{1,3})", text or "")
    r5yr = re.search(r"Risk Analysis: MPT Statistics \(5Yr\)\s+(\d{1,3})", text or "")

    perf_page = int(perf.group(1)) if perf else None
    sc_page   = int(sc.group(1))   if sc   else None
    fs_page   = int(fs.group(1))   if fs   else None
    cy_page   = int(cy.group(1))   if cy   else None
    r3yr_page = int(r3yr.group(1)) if r3yr else None
    r5yr_page = int(r5yr.group(1)) if r5yr else None

    st.subheader("Table of Contents Pages")
    st.write(f"- Fund Performance Current vs Proposed Comparison : {perf_page}")
    st.write(f"- Fund Performance Calendar Year : {cy_page}")
    st.write(f"- MPT 3Yr Risk Analysis : {r3yr_page}")
    st.write(f"- MPT 5Yr Risk Analysis : {r5yr_page}")
    st.write(f"- Fund Scorecard:   {sc_page}")
    st.write(f"- Fund Factsheets :  {fs_page}")
    


    # Store in session state for future reference
    st.session_state['performance_page'] = perf_page
    st.session_state['scorecard_page']   = sc_page
    st.session_state['factsheets_page']  = fs_page
    st.session_state['calendar_year_page'] = cy_page
    st.session_state['r3yr_page'] = r3yr_page
    st.session_state['r5yr_page'] = r5yr_page

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
import streamlit as st
import re
import pdfplumber
import pandas as pd
import yfinance as yf

def infer_fund_type_guess(ticker):
    """Infer 'Active' or 'Passive' based on Yahoo Finance info (name and summary)."""
    try:
        if not ticker:
            return ""
        info = yf.Ticker(ticker).info
        name = (info.get("longName") or info.get("shortName") or "").lower()
        summary = (info.get("longBusinessSummary") or "").lower()
        if "index" in name or "index" in summary:
            return "Passive"
        if "track" in summary and "index" in summary:
            return "Passive"
        if "actively managed" in summary or "actively-managed" in summary:
            return "Active"
        if "outperform" in summary or "manager selects" in summary:
            return "Active"
        return ""
    except Exception:
        return ""

def extract_scorecard_blocks(pdf, scorecard_page):
    metric_labels = [
        "Manager Tenure", "Excess Performance (3Yr)", "Excess Performance (5Yr)",
        "Peer Return Rank (3Yr)", "Peer Return Rank (5Yr)", "Expense Ratio Rank",
        "Sharpe Ratio Rank (3Yr)", "Sharpe Ratio Rank (5Yr)", "R-Squared (3Yr)",
        "R-Squared (5Yr)", "Sortino Ratio Rank (3Yr)", "Sortino Ratio Rank (5Yr)",
        "Tracking Error Rank (3Yr)", "Tracking Error Rank (5Yr)"
    ]
    pages, fund_blocks, fund_name, metrics = [], [], None, []
    for p in pdf.pages[scorecard_page-1:]:
        pages.append(p.extract_text() or "")
    lines = "\n".join(pages).splitlines()
    for line in lines:
        if not any(metric in line for metric in metric_labels) and line.strip():
            if fund_name and metrics:
                fund_blocks.append({"Fund Name": fund_name, "Metrics": metrics})
            fund_name = re.sub(r"Fund (Meets Watchlist Criteria|has been placed on watchlist for not meeting .* out of 14 criteria)", "", line.strip()).strip()
            metrics = []
        for metric in metric_labels:
            if metric in line:
                m = re.match(r"^(.*?)\s+(Pass|Review|Fail)\s*(.*)", line.strip())
                if m:
                    metric_name, status, info = m.groups()
                    metrics.append({"Metric": metric_name, "Status": status, "Info": info.strip()})
    if fund_name and metrics:
        fund_blocks.append({"Fund Name": fund_name, "Metrics": metrics})
    return fund_blocks

def extract_fund_tickers(pdf, performance_page, fund_names, factsheets_page=None):
    end_page = factsheets_page-1 if factsheets_page else len(pdf.pages)
    all_lines, perf_text = [], ""
    for p in pdf.pages[performance_page-1:end_page]:
        txt = p.extract_text() or ""
        perf_text += txt + "\n"
        all_lines.extend(txt.splitlines())
    mapping = {}
    for ln in all_lines:
        m = re.match(r"(.+?)\s+([A-Z]{1,5})$", ln.strip())
        if not m:
            continue
        raw_name, ticker = m.groups()
        norm = re.sub(r'[^A-Za-z0-9 ]+', '', raw_name).strip().lower()
        mapping[norm] = ticker
    tickers = {}
    for name in fund_names:
        norm_expected = re.sub(r'[^A-Za-z0-9 ]+', '', name).strip().lower()
        found = next((t for raw, t in mapping.items() if raw.startswith(norm_expected)), None)
        tickers[name] = found
    total = len(fund_names)
    found_count = sum(1 for t in tickers.values() if t)
    if found_count < total:
        all_tks = re.findall(r'\b([A-Z]{1,5})\b', perf_text)
        seen = []
        for tk in all_tks:
            if tk not in seen:
                seen.append(tk)
        tickers = {name: (seen[i] if i < len(seen) else "") for i, name in enumerate(fund_names)}
    return {k: (v if v else "") for k, v in tickers.items()}

def scorecard_to_ips(fund_blocks, fund_types, tickers):
    # Maps for converting scorecard to IPS criteria (from your logic)
    metrics_order = [
        "Manager Tenure", "Excess Performance (3Yr)", "Excess Performance (5Yr)",
        "Peer Return Rank (3Yr)", "Peer Return Rank (5Yr)", "Expense Ratio Rank",
        "Sharpe Ratio Rank (3Yr)", "Sharpe Ratio Rank (5Yr)", "R-Squared (3Yr)",
        "R-Squared (5Yr)", "Sortino Ratio Rank (3Yr)", "Sortino Ratio Rank (5Yr)",
        "Tracking Error Rank (3Yr)", "Tracking Error Rank (5Yr)",
    ]
    active_map  = [0,1,3,6,10,2,4,7,11,5,None]
    passive_map = [0,8,3,6,12,9,4,7,13,5,None]
    ips_labels = [f"IPS Investment Criteria {i+1}" for i in range(11)]
    ips_results, raw_results = [], []
    for fund in fund_blocks:
        fund_name = fund["Fund Name"]
        fund_type = fund_types.get(fund_name, "Passive" if "index" in fund_name.lower() else "Active")
        metrics = fund["Metrics"]
        scorecard_status = [next((m["Status"] for m in metrics if m["Metric"] == label), None) for label in metrics_order]
        idx_map = passive_map if fund_type == "Passive" else active_map
        ips_status = [scorecard_status[m_idx] if m_idx is not None else "Pass" for m_idx in idx_map]
        review_fail = sum(1 for status in ips_status if status in ["Review","Fail"])
        watch_status = "FW" if review_fail >= 6 else "IW" if review_fail >= 5 else "NW"
        def iconify(status): return "✔" if status == "Pass" else "✗" if status in ("Review", "Fail") else ""
        row = {
            "Fund Name": fund_name,
            "Ticker": tickers.get(fund_name, ""),
            "Fund Type": fund_type,
            **{ips_labels[i]: iconify(ips_status[i]) for i in range(11)},
            "IPS Watch Status": watch_status,
        }
        ips_results.append(row)
        raw_results.append({
            "Fund Name": fund_name,
            "Ticker": tickers.get(fund_name, ""),
            "Fund Type": fund_type,
            **{ips_labels[i]: ips_status[i] for i in range(11)},
            "IPS Watch Status": watch_status,
        })
    return pd.DataFrame(ips_results), pd.DataFrame(raw_results)

def watch_status_color(val):
    if val == "FW":
        return "background-color:#f8d7da; color:#c30000; font-weight:600;"
    if val == "IW":
        return "background-color:#fff3cd; color:#B87333; font-weight:600;"
    if val == "NW":
        return "background-color:#d6f5df; color:#217a3e; font-weight:600;"
    return ""

def step3_5_6_scorecard_and_ips(pdf, scorecard_page, performance_page, factsheets_page, total_options):
    # 1. Extract scorecard blocks
    fund_blocks = extract_scorecard_blocks(pdf, scorecard_page)
    fund_names = [fund["Fund Name"] for fund in fund_blocks]
    if not fund_blocks:
        st.error("Could not extract fund scorecard blocks. Check the PDF and page number.")
        return
    # 2. Extract tickers
    tickers = extract_fund_tickers(pdf, performance_page, fund_names, factsheets_page)
    # 3. Guess fund type using Yahoo (overridable)
    fund_type_guesses = []
    for name in fund_names:
        guess = ""
        if tickers.get(name):
            guess = infer_fund_type_guess(tickers[name])
        if guess == "Passive":
            fund_type_guesses.append("Passive")
        else:
            fund_type_guesses.append("Active")
    fund_type_defaults = ["Passive" if "index" in n.lower() else "Active" for n in fund_names]
    use_guess = st.checkbox(
        "Prefill Fund Type with Yahoo Finance guess instead of default (index = Passive, else Active)",
        value=True
    )
    prefill_fund_type = fund_type_guesses if use_guess else fund_type_defaults
    df_types = pd.DataFrame({
        "Fund Name": fund_names,
        "Ticker": [tickers[name] for name in fund_names],
        "Fund Type Guess": fund_type_guesses,
        "Fund Type": prefill_fund_type,
    })
    edited_types = st.data_editor(
        df_types,
        column_config={
            "Fund Type": st.column_config.SelectboxColumn("Fund Type", options=["Active", "Passive"]),
        },
        disabled=["Fund Name", "Ticker", "Fund Type Guess"],
        hide_index=True,
        key="data_editor_fundtype",
        use_container_width=True,
    )
    fund_types = {row["Fund Name"]: row["Fund Type"] for _, row in edited_types.iterrows()}
    # 4. Convert to IPS screening
    df_icon, df_raw = scorecard_to_ips(fund_blocks, fund_types, tickers)
    # 5. Show table (short labels)
    if not df_icon.empty:
        display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
        display_df = df_icon.rename(columns=display_columns)
        st.markdown(
            '<div class="watch-key" style="margin-bottom: 1em;">'
            '<span style="background:#d6f5df; color:#217a3e; padding:0.07em 0.55em; border-radius:2px;">NW</span> '
            '(No Watch) &nbsp;'
            '<span style="background:#fff3cd; color:#B87333; padding:0.07em 0.55em; border-radius:2px;">IW</span> '
            '(Informal Watch) &nbsp;'
            '<span style="background:#f8d7da; color:#c30000; padding:0.07em 0.55em; border-radius:2px;">FW</span> '
            '(Formal Watch)</div>', unsafe_allow_html=True
        )
        styled = display_df.style.applymap(watch_status_color, subset=["IPS Watch Status"])
        st.dataframe(styled, use_container_width=True, hide_index=True)
        st.download_button(
            "Download CSV",
            data=df_raw.to_csv(index=False),
            file_name="ips_screening_table.csv",
            mime="text/csv",
        )
    else:
        st.info("No IPS screening data available.")
    # 6. Save to session state for downstream steps
    st.session_state["fund_blocks"] = fund_blocks
    st.session_state["fund_types"] = fund_types
    st.session_state["fund_tickers"] = tickers
    st.session_state["ips_icon_table"] = df_icon
    st.session_state["ips_raw_table"] = df_raw

    # 7. Extract & save fund performance data for later steps
    perf_data = extract_performance_table(
        pdf,
        performance_page,
        fund_names,
        factsheets_page
    )
    # Attach tickers (already extracted earlier!)
    for itm in perf_data:
        itm["Ticker"] = tickers.get(itm["Fund Scorecard Name"], "")

    # Save for Step 7 and others
    st.session_state["fund_performance_data"] = perf_data
    st.session_state["tickers"] = tickers  # Keep ticker mapping for legacy steps

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
# === Step 6: Fund Factsheets ===
def step6_process_factsheets(pdf, fund_names):
    st.subheader("Step 6: Fund Factsheets Section")
    factsheet_start = st.session_state.get("factsheets_page")
    total_declared = st.session_state.get("total_options")
    performance_data = [
        {"Fund Scorecard Name": name, "Ticker": ticker}
        for name, ticker in st.session_state.get("tickers", {}).items()
    ]

    if not factsheet_start:
        st.error("❌ 'Fund Factsheets' page number not found in TOC.")
        return

    matched_factsheets = []
    # Iterate pages from factsheet_start to end
    for i in range(factsheet_start - 1, len(pdf.pages)):
        page = pdf.pages[i]
        words = page.extract_words(use_text_flow=True)
        header_words = [w['text'] for w in words if w['top'] < 100]
        first_line = " ".join(header_words).strip()

        if not first_line or "Benchmark:" not in first_line or "Expense Ratio:" not in first_line:
            continue

        ticker_match = re.search(r"\b([A-Z]{5})\b", first_line)
        ticker = ticker_match.group(1) if ticker_match else ""
        fund_name_raw = first_line.split(ticker)[0].strip() if ticker else first_line

        best_score = 0
        matched_name = matched_ticker = ""
        for item in performance_data:
            ref = f"{item['Fund Scorecard Name']} {item['Ticker']}".strip()
            score = fuzz.token_sort_ratio(f"{fund_name_raw} {ticker}".lower(), ref.lower())
            if score > best_score:
                best_score, matched_name, matched_ticker = score, item['Fund Scorecard Name'], item['Ticker']

        def extract_field(label, text, stop=None):
            try:
                start = text.index(label) + len(label)
                rest = text[start:]
                if stop and stop in rest:
                    return rest[:rest.index(stop)].strip()
                return rest.split()[0]
            except Exception:
                return ""

        benchmark = extract_field("Benchmark:", first_line, "Category:")
        category  = extract_field("Category:", first_line, "Net Assets:")
        net_assets= extract_field("Net Assets:", first_line, "Manager Name:")
        manager   = extract_field("Manager Name:", first_line, "Avg. Market Cap:")
        avg_cap   = extract_field("Avg. Market Cap:", first_line, "Expense Ratio:")
        expense   = extract_field("Expense Ratio:", first_line)

        matched_factsheets.append({
            "Page #": i + 1,
            "Parsed Fund Name": fund_name_raw,
            "Parsed Ticker": ticker,
            "Matched Fund Name": matched_name,
            "Matched Ticker": matched_ticker,
            "Benchmark": benchmark,
            "Category": category,
            "Net Assets": net_assets,
            "Manager Name": manager,
            "Avg. Market Cap": avg_cap,
            "Expense Ratio": expense,
            "Match Score": best_score,
            "Matched": "✅" if best_score > 20 else "❌"
        })

    df_facts = pd.DataFrame(matched_factsheets)
    st.session_state['fund_factsheets_data'] = matched_factsheets

    display_df = df_facts[[
        "Matched Fund Name", "Matched Ticker", "Benchmark", "Category",
        "Net Assets", "Manager Name", "Avg. Market Cap", "Expense Ratio", "Matched"
    ]].rename(columns={"Matched Fund Name": "Fund Name", "Matched Ticker": "Ticker"})

    st.dataframe(display_df, use_container_width=True)

    matched_count = sum(1 for r in matched_factsheets if r["Matched"] == "✅")
    if not st.session_state.get("suppress_matching_confirmation", False):
        st.write(f"Matched {matched_count} of {len(matched_factsheets)} factsheet pages.")
        if matched_count == total_declared:
            st.success(f"All {matched_count} funds matched the declared Total Options from Page 1.")
        else:
            st.error(f"Mismatch: Page 1 declared {total_declared}, but only matched {matched_count}.")

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# === Step 7: QTD / 1Yr / 3Yr / 5Yr / 10Yr / Net Expense Ratio & Bench QTD ===
def step7_extract_returns(pdf):
    import re
    import pandas as pd
    import streamlit as st
    from rapidfuzz import fuzz

    st.subheader("Step 7: QTD / 1Yr / 3Yr / 5Yr / 10Yr / Net Expense & Benchmark QTD")

    # 1) Where to scan
    perf_page = st.session_state.get("performance_page")
    end_page  = st.session_state.get("calendar_year_page") or (len(pdf.pages) + 1)
    perf_data = st.session_state.get("fund_performance_data", [])
    if perf_page is None or not perf_data:
        st.error("❌ Run Step 5 first to populate performance data.")
        return

    # 2) Prep output slots
    fields = ["QTD", "1Yr", "3Yr", "5Yr", "10Yr", "Net Expense Ratio", "Bench QTD", "Bench 3Yr", "Bench 5Yr"]
    for itm in perf_data:
        for f in fields:
            itm.setdefault(f, None)

    # 3) Gather every nonblank line in the Performance section
    lines = []
    for pnum in range(perf_page - 1, end_page - 1):
        txt = pdf.pages[pnum].extract_text() or ""
        lines += [ln.strip() for ln in txt.splitlines() if ln.strip()]

    # 4) Regex to pull decimal tokens (with optional % and parentheses)
    num_rx = re.compile(r"\(?-?\d+\.\d+%?\)?")

    matched = 0
    for item in perf_data:
        name = item["Fund Scorecard Name"]
        tk   = item["Ticker"].upper().strip()

        # a) Exact-ticker match
        idx = next(
            (i for i, ln in enumerate(lines)
             if re.search(rf"\b{re.escape(tk)}\b", ln)),
            None
        )
        # b) Fuzzy-name fallback
        if idx is None:
            scores = [(i, fuzz.token_sort_ratio(name.lower(), ln.lower()))
                      for i, ln in enumerate(lines)]
            best_i, best_score = max(scores, key=lambda x: x[1])
            if best_score > 60:
                idx = best_i
            else:
                st.warning(f"⚠️ {name} ({tk}): no match found.")
                continue

        # c) Pull fund numbers from line above (and two above if needed)
        raw = num_rx.findall(lines[idx - 1]) if idx >= 1 else []
        if len(raw) < 8 and idx >= 2:
            raw = num_rx.findall(lines[idx - 2]) + raw
        clean = [n.strip("()%").rstrip("%") for n in raw]
        if len(clean) < 8:
            clean += [None] * (8 - len(clean))

        # d) Map fund returns & net expense
        item["QTD"]               = clean[0]
        item["1Yr"]               = clean[2]
        item["3Yr"]               = clean[3]
        item["5Yr"]               = clean[4]
        item["10Yr"]              = clean[5]
        item["Net Expense Ratio"] = clean[-2]

        # e) Pull benchmark QTD, 3Yr, and 5Yr from the very next line (or one more down)
        bench_raw = []
        if idx + 1 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 1])
        if len(bench_raw) < 1 and idx + 2 < len(lines):
            bench_raw = num_rx.findall(lines[idx + 2])
        bench_clean = [n.strip("()%").rstrip("%") for n in bench_raw]

        item["Bench QTD"] = bench_clean[0] if bench_clean else None
        item["Bench 3Yr"] = bench_clean[3] if len(bench_clean) > 3 else None
        item["Bench 5Yr"] = bench_clean[4] if len(bench_clean) > 4 else None

        matched += 1

    # 5) Save & display
    st.session_state["fund_performance_data"] = perf_data
    df = pd.DataFrame(perf_data)

    st.success(f"✅ Matched {matched} fund(s) with return data.")
    for itm in perf_data:
        missing = [f for f in fields if not itm.get(f)]
        if missing:
            st.warning(f"⚠️ Incomplete for {itm['Fund Scorecard Name']} ({itm['Ticker']}): missing {', '.join(missing)}")

    st.dataframe(
        df[["Fund Scorecard Name", "Ticker"] + fields],
        use_container_width=True
    )

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# === Step 8 Calendar Year Returns (funds + benchmarks) ===
def step8_calendar_returns(pdf):
    import re, streamlit as st, pandas as pd

    st.subheader("Step 8: Calendar Year Returns")

    # 1) Figure out section bounds
    cy_page  = st.session_state.get("calendar_year_page")
    end_page = st.session_state.get("r3yr_page", len(pdf.pages) + 1)
    if cy_page is None:
        st.error("❌ 'Fund Performance: Calendar Year' not found in TOC.")
        return

    # 2) Pull every line from that section
    all_lines = []
    for p in pdf.pages[cy_page-1 : end_page-1]:
        all_lines.extend((p.extract_text() or "").splitlines())

    # 3) Identify header & years
    header = next((ln for ln in all_lines if "Ticker" in ln and re.search(r"20\d{2}", ln)), None)
    if not header:
        st.error("❌ Couldn’t find header row with 'Ticker' + year.")
        return
    years = re.findall(r"\b20\d{2}\b", header)
    num_rx = re.compile(r"-?\d+\.\d+%?")

    # — A) Funds themselves —
    fund_map     = st.session_state.get("tickers", {})
    fund_records = []
    for name, tk in fund_map.items():
        ticker = (tk or "").upper()
        idx    = next((i for i, ln in enumerate(all_lines) if ticker in ln.split()), None)
        raw    = num_rx.findall(all_lines[idx-1]) if idx not in (None, 0) else []
        vals   = raw[:len(years)] + [None] * (len(years) - len(raw))
        rec    = {"Name": name, "Ticker": ticker}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        fund_records.append(rec)

    df_fund = pd.DataFrame(fund_records)
    if not df_fund.empty:
        st.markdown("**Fund Calendar‑Year Returns**")
        st.dataframe(df_fund[["Name", "Ticker"] + years], use_container_width=True)
        st.session_state["step8_returns"] = fund_records

    # — B) Benchmarks matched back to each fund’s ticker —
    facts         = st.session_state.get("fund_factsheets_data", [])
    bench_records = []
    for f in facts:
        bench_name = f.get("Benchmark", "").strip()
        fund_tkr   = f.get("Matched Ticker", "")
        if not bench_name:
            continue

        # find the first line containing the benchmark name
        idx = next((i for i, ln in enumerate(all_lines) if bench_name in ln), None)
        if idx is None:
            continue
        raw  = num_rx.findall(all_lines[idx])
        vals = raw[:len(years)] + [None] * (len(years) - len(raw))
        rec  = {"Name": bench_name, "Ticker": fund_tkr}
        rec.update({years[i]: vals[i] for i in range(len(years))})
        bench_records.append(rec)

    df_bench = pd.DataFrame(bench_records)
    if not df_bench.empty:
        st.markdown("**Benchmark Calendar‑Year Returns**")
        st.dataframe(df_bench[["Name", "Ticker"] + years], use_container_width=True)
        st.session_state["benchmark_calendar_year_returns"] = bench_records
    else:
        st.warning("No benchmark returns extracted.")

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# === Step 9: 3‑Yr Risk Analysis – Match & Extract MPT Stats (hidden matching) ===
def step9_risk_analysis_3yr(pdf):
    import re, streamlit as st, pandas as pd
    from rapidfuzz import fuzz

    st.subheader("Step 9: Risk Analysis (3Yr) – MPT Statistics")

    # 1) Get your fund→ticker map
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # 2) Locate the “Risk Analysis: MPT Statistics (3Yr)” page
    start_page = st.session_state.get("r3yr_page")
    if not start_page:
        st.error("❌ ‘Risk Analysis: MPT Statistics (3Yr)’ page not found; run Step 2 first.")
        return

    # 3) Scan forward until you’ve seen each ticker (no display)
    locs = {}
    for pnum in range(start_page, len(pdf.pages) + 1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            tokens = ln.split()
            for fname, tk in fund_map.items():
                if fname in locs: 
                    continue
                if tk.upper() in tokens:
                    locs[fname] = {"page": pnum, "line": li}
        if len(locs) == len(fund_map):
            break

    # 4) Extract the first four numeric MPT stats from that same line
    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, info in locs.items():
        page = pdf.pages[info["page"]-1]
        lines = (page.extract_text() or "").splitlines()
        line = lines[info["line"]]
        nums = num_rx.findall(line)
        nums += [None] * (4 - len(nums))
        alpha, beta, up, down = nums[:4]
        results.append({
            "Fund Name":               name,
            "Ticker":                  fund_map[name].upper(),
            "3 Year Alpha":            alpha,
            "3 Year Beta":             beta,
            "3 Year Upside Capture":   up,
            "3 Year Downside Capture": down
        })

    # 5) Display final table only
    st.session_state["step9_mpt_stats"] = results
    df = pd.DataFrame(results)
    st.dataframe(df, use_container_width=True)


#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# === Step 10: Risk Analysis (5Yr) – Match & Extract MPT Statistics ===
def step10_risk_analysis_5yr(pdf):
    import re, streamlit as st, pandas as pd

    st.subheader("Step 10: Risk Analysis (5Yr) – MPT Statistics")

    # 1) Your fund→ticker map from Step 5
    fund_map = st.session_state.get("tickers", {})
    if not fund_map:
        st.error("❌ No ticker mapping found. Run Step 5 first.")
        return

    # 2) Locate the “Risk Analysis: MPT Statistics (5Yr)” section
    section_page = next(
        (i for i, pg in enumerate(pdf.pages, start=1)
         if "Risk Analysis: MPT Statistics (5Yr)" in (pg.extract_text() or "")),
        None
    )
    if section_page is None:
        st.error("❌ Could not find ‘Risk Analysis: MPT Statistics (5Yr)’ section.")
        return

    # 3) Under‑the‑hood: scan pages until each ticker is located
    locs = {}
    total = len(fund_map)
    for pnum in range(section_page, len(pdf.pages) + 1):
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()
        for li, ln in enumerate(lines):
            tokens = ln.split()
            for name, tk in fund_map.items():
                if name in locs:
                    continue
                if tk.upper() in tokens:
                    locs[name] = {"page": pnum, "line": li}
        if len(locs) == total:
            break

    # 4) Wrap‑aware extraction of the first four floats after each ticker line
    num_rx = re.compile(r"-?\d+\.\d+")
    results = []
    for name, tk in fund_map.items():
        info = locs.get(name)
        vals = [None] * 4
        if info:
            page = pdf.pages[info["page"] - 1]
            text_lines = (page.extract_text() or "").splitlines()
            idx = info["line"]
            nums = []
            # look on the line of the ticker and up to the next 2 lines
            for j in range(idx, min(idx + 3, len(text_lines))):
                nums += num_rx.findall(text_lines[j])
                if len(nums) >= 4:
                    break
            nums += [None] * (4 - len(nums))
            vals = nums[:4]
        else:
            st.warning(f"⚠️ {name} ({tk.upper()}): not found after page {section_page}.")

        alpha5, beta5, up5, down5 = vals
        results.append({
            "Fund Name":               name,
            "Ticker":                  tk.upper(),
            "5 Year Alpha":            alpha5,
            "5 Year Beta":             beta5,
            "5 Year Upside Capture":   up5,
            "5 Year Downside Capture": down5,
        })

    # 5) Save & display only the consolidated table
    st.session_state["step10_mpt_stats"] = results
    df = pd.DataFrame(results)
    st.dataframe(df, use_container_width=True)


#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# === Step 11: Combined MPT Statistics Summary ===
def step11_create_summary(pdf=None):
    import pandas as pd
    import streamlit as st

    st.subheader("Step 11: MPT Statistics Summary")

    # 1) Load your 3‑Yr and 5‑Yr stats from session state
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    mpt5 = st.session_state.get("step10_mpt_stats", [])
    if not mpt3 or not mpt5:
        st.error("❌ Missing MPT stats. Run Steps 9 & 10 first.")
        return

    # 2) Build DataFrames
    df3 = pd.DataFrame(mpt3)  # contains "3 Year Alpha", "3 Year Beta", etc.
    df5 = pd.DataFrame(mpt5)  # contains "5 Year Alpha", "5 Year Beta", etc.

    # 3) Merge on Fund Name & Ticker
    df = pd.merge(
        df3,
        df5,
        on=["Fund Name", "Ticker"],
        how="outer",
        suffixes=("_3yr", "_5yr")
    )

    # 4) Build the Investment Manager column
    df.insert(0, "Investment Manager", df["Fund Name"] + " (" + df["Ticker"] + ")")

    # 5) Select & order the columns
    df = df[[
        "Investment Manager",
        "3 Year Alpha",
        "5 Year Alpha",
        "3 Year Beta",
        "5 Year Beta",
        "3 Year Upside Capture",
        "3 Year Downside Capture",
        "5 Year Upside Capture",
        "5 Year Downside Capture"
    ]]

    # 6) Display
    st.session_state["step11_summary"] = df.to_dict("records")
    st.dataframe(df)

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# === Step 12: Extract “FUND FACTS” & Its Table Details in One Go ===
def step12_process_fund_facts(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.subheader("Step 12: Fund Facts")

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # the exact labels and the order you want them in the table
    labels = [
        "Manager Tenure Yrs.",
        "Expense Ratio",
        "Expense Ratio Rank",
        "Total Number of Holdings",
        "Turnover Ratio"
    ]

    records = []
    # scan each factsheet page
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = pdf.pages[pnum-1].extract_text().splitlines()

        for idx, line in enumerate(lines):
            if line.lstrip().upper().startswith("FUND FACTS"):
                # grab the next 8 lines (should contain your 5 labels)
                snippet = lines[idx+1 : idx+1+8]
                rec = {"Fund Name": fund_name, "Ticker": ticker}
                for lab in labels:
                    val = None
                    for ln in snippet:
                        norm = " ".join(ln.strip().split())
                        if norm.startswith(lab):
                            rest = norm[len(lab):].strip(" :\t")
                            m = re.match(r"(-?\d+\.\d+)", rest)
                            val = m.group(1) if m else (rest.split()[0] if rest else None)
                            break
                    rec[lab] = val
                records.append(rec)
                break  # move on to the next page once Fund Facts is processed

    if not records:
        st.warning("No Fund Facts tables found.")
        return

    # save & show
    st.session_state["step12_fund_facts_table"] = records
    df = pd.DataFrame(records)
    st.dataframe(df, use_container_width=True)

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# === Step 13: Extract Risk‑Adjusted Returns Metrics ===
def step13_process_risk_adjusted_returns(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.subheader("Step 13: Risk‑Adjusted Returns")

    fs_start   = st.session_state.get("factsheets_page")
    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not fs_start or not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    # map factsheet pages to fund name & ticker
    page_map = {
        f["Page #"]: (f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    # which metrics to pull
    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    num_rx  = re.compile(r"-?\d+\.\d+")

    records = []
    for pnum in range(fs_start, len(pdf.pages) + 1):
        if pnum not in page_map:
            continue
        fund_name, ticker = page_map[pnum]
        lines = (pdf.pages[pnum-1].extract_text() or "").splitlines()

        # find the heading
        for idx, line in enumerate(lines):
            norm = " ".join(line.strip().split()).upper()
            if norm.startswith("RISK-ADJUSTED RETURNS"):
                snippet = lines[idx+1 : idx+1+6]  # grab next few lines
                rec = {"Fund Name": fund_name, "Ticker": ticker}

                for metric in metrics:
                    # find the snippet line for this metric
                    text_line = next(
                        ( " ".join(ln.strip().split())
                          for ln in snippet
                          if ln.strip().upper().startswith(metric.upper()) ),
                        None
                    ) or ""
                    # extract up to 4 numbers
                    nums = num_rx.findall(text_line)
                    nums += [None] * (4 - len(nums))

                    # assign into rec
                    rec[f"{metric} 1Yr"]  = nums[0]
                    rec[f"{metric} 3Yr"]  = nums[1]
                    rec[f"{metric} 5Yr"]  = nums[2]
                    rec[f"{metric} 10Yr"] = nums[3]

                records.append(rec)
                break  # done with this page

    if not records:
        st.warning("No 'RISK‑ADJUSTED RETURNS' tables found.")
        return

    # save & show
    st.session_state["step13_risk_adjusted_table"] = records
    df = pd.DataFrame(records)
    st.dataframe(df, use_container_width=True)

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────

# == Step 14: Peer Risk-Adjusted Return Rank ==
def step14_extract_peer_risk_adjusted_return_rank(pdf):
    import re
    import streamlit as st
    import pandas as pd

    st.subheader("Step 14: Peer Risk-Adjusted Return Rank")

    factsheets = st.session_state.get("fund_factsheets_data", [])
    if not factsheets:
        st.error("❌ Run Step 6 first to populate your factsheet pages.")
        return

    page_map = {
        f["Page #"]:(f["Matched Fund Name"], f["Matched Ticker"])
        for f in factsheets
    }

    metrics = ["Sharpe Ratio", "Information Ratio", "Sortino Ratio"]
    records = []

    for pnum, (fund, ticker) in page_map.items():
        page = pdf.pages[pnum-1]
        text = page.extract_text() or ""
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]

        # 1) locate the Risk-Adjusted Returns header
        try:
            risk_idx = next(i for i, ln in enumerate(lines)
                            if "RISK-ADJUSTED RETURNS" in ln.upper())
        except StopIteration:
            st.warning(f"⚠️ {fund} ({ticker}): Risk-Adjusted header not found.")
            continue

        # 2) find all “1 Yr 3 Yrs 5 Yrs 10 Yrs” lines after that
        header_idxs = [i for i, ln in enumerate(lines)
                       if re.match(r"1\s*Yr", ln)]
        peer_header_idxs = [i for i in header_idxs if i > risk_idx]

        if not peer_header_idxs:
            st.warning(f"⚠️ {fund} ({ticker}): peer header not found.")
            continue

        # take the *second* header occurrence (first is Risk-Adjusted, next is Peer)
        peer_hdr = peer_header_idxs[0] if len(peer_header_idxs)==1 else peer_header_idxs[1]

        rec = {"Fund Name": fund, "Ticker": ticker}

        # 3) read the three lines immediately below that header
        for offset, metric in enumerate(metrics, start=1):
            if peer_hdr + offset < len(lines):
                parts = lines[peer_hdr + offset].split()
                # parts[0:2] = metric name words, parts[2:6] = the four integer ranks
                vals = parts[2:6] if len(parts) >= 6 else []
            else:
                vals = []

            # fill into record (pad with None if missing)
            for idx, period in enumerate(["1Yr","3Yr","5Yr","10Yr"]):
                rec[f"{metric} {period}"] = vals[idx] if idx < len(vals) else None

            if len(vals) < 4:
                st.warning(f"⚠️ {fund} ({ticker}): only {len(vals)} peer values found for '{metric}'.")

        records.append(rec)

    if not records:
        st.warning("❌ No Peer Risk-Adjusted Return Rank data extracted.")
        return

    df = pd.DataFrame(records)
    st.session_state["step14_peer_rank_table"] = records
    st.dataframe(df, use_container_width=True)


#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
#Step 15
def step15_display_selected_fund():
    import pandas as pd
    import streamlit as st
    import re

    st.subheader("Step 15: Single Fund Details")
    facts = st.session_state.get("fund_factsheets_data", [])
    if not facts:
        st.info("Run Steps 1–14 to populate data before viewing fund details.")
        return

    fund_names = [f["Matched Fund Name"] for f in facts]
    selected_fund = st.selectbox("Select a fund to view details:", fund_names)
    st.session_state.selected_fund = selected_fund

    st.write(f"Details for: {selected_fund}")

    # --- Slide 1 Table: IPS Results ---
    ips_icon_table = st.session_state.get("ips_icon_table")
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected_fund), None)

    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        if not row.empty:
            row_dict = row.iloc[0].to_dict()
            display_columns = {f"IPS Investment Criteria {i+1}": str(i+1) for i in range(11)}
            row_df = pd.DataFrame([{
                "Category": fs_rec.get("Category", "") if fs_rec else "",
                "Time Period": st.session_state.get("report_date", ""),
                "Plan Assets": "$",  # Or replace with actual variable if you store this elsewhere!
                **{display_columns.get(k, k): v for k, v in row_dict.items() if k.startswith("IPS Investment Criteria")},
                "IPS Status": row_dict.get("IPS Watch Status", "")
            }])

            def color_bool(v): return "background-color: green" if v == "✔" else ("background-color: red" if v == "✗" else "")
            def style_status(v):
                if v == "NW": return "background-color: green; color: white"
                if v == "IW": return "background-color: orange; color: white"
                if v == "FW": return "background-color: red; color: white"
                return ""
            styled = row_df.style.applymap(color_bool, subset=[str(i) for i in range(1, 12)]).applymap(style_status, subset=["IPS Status"])
            st.dataframe(styled, use_container_width=True)
        else:
            st.warning("No IPS screening result found for selected fund.")
    else:
        st.warning("IPS screening table not found. Run earlier steps first.")

    # --- Slide 2 Table 1 ---
    st.markdown("**Slide 2 Table 1**")
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == selected_fund), {})
    inv_mgr = f"{selected_fund} ({perf_item.get('Ticker','')})"
    net_exp = perf_item.get("Net Expense Ratio", "")
    if net_exp and not str(net_exp).endswith("%"):
        net_exp = f"{net_exp}%"
    df_slide2 = pd.DataFrame([{
        "Investment Manager": inv_mgr,
        "Net Expense Ratio":  net_exp
    }])
    st.dataframe(df_slide2, use_container_width=True)

    # --- Slide 2 Table 2 ---
    st.markdown("**Slide 2 Table 2**")
    date_label = st.session_state.get("report_date", "QTD")
    def append_pct(val):
        s = str(val) if val is not None else ""
        return s if s.endswith("%") or s=="" else f"{s}%"
    qtd   = append_pct(perf_item.get("QTD",""))
    one   = append_pct(perf_item.get("1Yr",""))
    three = append_pct(perf_item.get("3Yr",""))
    five  = append_pct(perf_item.get("5Yr",""))
    ten   = append_pct(perf_item.get("10Yr",""))
    row = {
        "Investment Manager": inv_mgr,
        date_label:           qtd,
        "1 Year":             one,
        "3 Year":             three,
        "5 Year":             five,
        "10 Year":            ten
    }
    df_slide2_2 = pd.DataFrame([row])
    st.dataframe(df_slide2_2, use_container_width=True)

    # --- Slide 2 Table 3 ---
    st.markdown("**Slide 2 Table 3**")
    fund_cy = st.session_state.get("step8_returns", [])
    bench_cy = st.session_state.get("benchmark_calendar_year_returns", [])
    if not fund_cy or not bench_cy:
        st.error("❌ No calendar year returns data found. Ensure Step 8 has been run correctly.")
        return
    fund_rec = next((r for r in fund_cy if r.get("Name") == selected_fund), None)
    if not fund_rec:
        st.error(f"❌ Could not find data for selected fund: {selected_fund}")
        return
    benchmark_name = selected_fund
    bench_rec = next((r for r in bench_cy if r.get("Name") == benchmark_name or r.get("Ticker") == fund_rec.get("Ticker")), None)
    if not bench_rec:
        st.error(f"❌ Could not find benchmark data for selected fund: {selected_fund}")
        return
    year_cols = [col for col in fund_rec.keys() if re.match(r"20\d{2}", col)]
    rows = []
    row_fund = {"Investment Manager": f"{selected_fund} ({fund_rec.get('Ticker','')})"}
    for year in year_cols:
        row_fund[year] = fund_rec.get(year, "")
    rows.append(row_fund)
    row_benchmark = {"Investment Manager": f"{bench_rec.get('Name', 'Benchmark')} ({bench_rec.get('Ticker', '')})"}
    for year in year_cols:
        row_benchmark[year] = bench_rec.get(year, "")
    rows.append(row_benchmark)
    df_slide2_3 = pd.DataFrame(rows, columns=["Investment Manager"] + year_cols)
    st.dataframe(df_slide2_3, use_container_width=True)

    # --- Slide 3 Table 1 ---
    st.markdown("**Slide 3 Table 1**")
    mpt3 = st.session_state.get("step9_mpt_stats", [])
    stats3 = next((r for r in mpt3 if r["Fund Name"] == selected_fund), {})
    mpt5 = st.session_state.get("step10_mpt_stats", [])
    stats5 = next((r for r in mpt5 if r["Fund Name"] == selected_fund), {})
    ticker = stats3.get("Ticker", stats5.get("Ticker", ""))
    inv_mgr = f"{selected_fund} ({ticker})"
    row = {
        "Investment Manager":        inv_mgr,
        "3 Year Alpha":              stats3.get("3 Year Alpha", ""),
        "5 Year Alpha":              stats5.get("5 Year Alpha", ""),
        "3 Year Beta":               stats3.get("3 Year Beta", ""),
        "5 Year Beta":               stats5.get("5 Year Beta", ""),
        "3 Year Upside Capture":     stats3.get("3 Year Upside Capture", ""),
        "3 Year Downside Capture":   stats3.get("3 Year Downside Capture", ""),
        "5 Year Upside Capture":     stats5.get("5 Year Upside Capture", ""),
        "5 Year Downside Capture":   stats5.get("5 Year Downside Capture", "")
    }
    df_slide3_1 = pd.DataFrame([row])
    st.dataframe(df_slide3_1, use_container_width=True)

    # --- Slide 3 Table 2 ---
    st.markdown("**Slide 3 Table 2**")
    risk_table = st.session_state.get("step13_risk_adjusted_table", [])
    peer_table = st.session_state.get("step14_peer_rank_table", [])
    risk_rec = next((r for r in risk_table if r["Fund Name"] == selected_fund), {})
    peer_rec = next((r for r in peer_table if r["Fund Name"] == selected_fund), {})
    ticker = risk_rec.get("Ticker") or peer_rec.get("Ticker", "")
    inv_mgr = f"{selected_fund} ({ticker})"
    def frac(metric, period):
        r = risk_rec.get(f"{metric} {period}", "")
        p = peer_rec.get(f"{metric} {period}", "")
        return f"{r} / {p}"
    row = {
        "Investment Manager": inv_mgr,
        "3 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "3Yr"),
        "5 Year Sharpe Ratio / Peer Ranking %": frac("Sharpe Ratio", "5Yr"),
        "3 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "3Yr"),
        "5 Year Sortino Ratio / Peer Ranking %": frac("Sortino Ratio", "5Yr"),
        "3 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "3Yr"),
        "5 Year Information Ratio / Peer Ranking %": frac("Information Ratio", "5Yr"),
    }
    df_slide3_2 = pd.DataFrame([row])
    st.dataframe(df_slide3_2, use_container_width=True)

    # --- Slide 4 Table 1 ---
    st.markdown("**Slide 4 Table 1**")
    blocks      = st.session_state.get("fund_blocks", [])
    block       = next((b for b in blocks if b["Fund Name"] == selected_fund), {})
    raw_tenure  = next((m["Info"] for m in block.get("Metrics", []) if m["Metric"] == "Manager Tenure"), "")
    m = re.search(r"(\d+(\.\d+)?)", raw_tenure)
    tenure = f"{m.group(1)} years" if m else raw_tenure
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == selected_fund), {})
    inv_mgr   = f"{selected_fund} ({perf_item.get('Ticker','')})"
    df_slide4 = pd.DataFrame([{
        "Investment Manager": inv_mgr,
        "Manager Tenure":     tenure
    }])
    st.dataframe(df_slide4, use_container_width=True)

    # --- Slide 4 Table 2 ---
    st.markdown("**Slide 4 Table 2**")
    facts = st.session_state.get("fund_factsheets_data", [])
    fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected_fund), None)
    perf_data = st.session_state.get("fund_performance_data", [])
    perf_item = next((p for p in perf_data if p["Fund Scorecard Name"] == selected_fund), None)
    inv_mgr    = f"{selected_fund} ({perf_item.get('Ticker','') if perf_item else ''})"
    assets     = fs_rec.get("Net Assets", "") if fs_rec else ""
    avg_cap    = fs_rec.get("Avg. Market Cap", "") if fs_rec else ""
    df_slide4_2 = pd.DataFrame([{
        "Investment Manager":             inv_mgr,
        "Assets Under Management":        assets,
        "Average Market Capitalization":  avg_cap
    }])
    st.dataframe(df_slide4_2, use_container_width=True)


def step16_bullet_points():
    import streamlit as st

    st.subheader("Step 16: Bullet Points")

    selected_fund = st.session_state.get("selected_fund")
    if not selected_fund:
        st.error("❌ No fund selected. Please select a fund in Step 15.")
        return

    perf_data = st.session_state.get("fund_performance_data", [])
    item = next((x for x in perf_data if x["Fund Scorecard Name"] == selected_fund), None)
    if not item:
        st.error(f"❌ Performance data for '{selected_fund}' not found.")
        return

    # Bullet 1: Performance vs. Benchmark, using template
    template = st.session_state.get("bullet_point_templates", [""])[0]
    b1 = template
    for fld, val in item.items():
        b1 = b1.replace(f"[{fld}]", str(val))
    st.markdown(f"- {b1}")

    # Get IPS status from icon table (guaranteed to match Slide 1 Table)
    ips_icon_table = st.session_state.get("ips_icon_table")
    ips_status = None
    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected_fund]
        ips_status = row.iloc[0]["IPS Watch Status"] if not row.empty else None

    # Bullet 2: Watch status and return comparison
    if ips_status == "NW":
        st.markdown("- This fund is **not on watch**.")
    else:
        status_label = (
            "Formal Watch" if ips_status == "FW" else
            "Informal Watch" if ips_status == "IW" else
            ips_status or "on watch"
        )

        three   = float(item.get("3Yr") or 0)
        bench3  = float(item.get("Bench 3Yr") or 0)
        five    = float(item.get("5Yr") or 0)
        bench5  = float(item.get("Bench 5Yr") or 0)
        bps3 = round((three - bench3) * 100, 1)
        bps5 = round((five  - bench5) * 100, 1)

        # Peer rank logic (safe handling)
        peer = st.session_state.get("step14_peer_rank_table", [])
        raw3 = next((r.get("Sharpe Ratio 3Yr") or r.get("Sharpe Ratio Rank 3Yr") for r in peer
                     if r.get("Fund Name") == selected_fund), None)
        raw5 = next((r.get("Sharpe Ratio 5Yr") or r.get("Sharpe Ratio Rank 5Yr") for r in peer
                     if r.get("Fund Name") == selected_fund), None)
        try:
            pos3 = "top" if raw3 and int(raw3) <= 50 else "bottom"
        except:
            pos3 = "bottom"
        try:
            pos5 = "top" if raw5 and int(raw5) <= 50 else "bottom"
        except:
            pos5 = "bottom"

        st.markdown(
            f"- The fund is now on **{status_label}**. Its 3‑year return trails the benchmark by "
            f"{bps3} bps ({three:.2f}% vs. {bench3:.2f}%) and its 5‑year return trails by "
            f"{bps5} bps ({five:.2f}% vs. {bench5:.2f}%). Its 3‑Yr Sharpe ranks in the {pos3} half of peers "
            f"and its 5‑Yr Sharpe ranks in the {pos5} half."
        )

    # Bullet 3: Action for Formal Watch only
    if ips_status == "FW":
        st.markdown("- **Action:** Consider replacing this fund.")

#── Build Powerpoint───────────────────────────────────────────────────────────────────────────────────────────────────────────────────

def step17_export_to_ppt_headings():
    import streamlit as st
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt
    from pptx.oxml.xmlchemy import OxmlElement
    from io import BytesIO
    import re
    import pandas as pd

    st.subheader("Step 17: Export to PowerPoint")

    # --- 1) Selected fund & category
    selected = st.session_state.get("selected_fund")
    if not selected:
        st.error("❌ No fund selected. Please select a fund in Step 15.")
        return

    facts = st.session_state.get("fund_factsheets_data", [])
    rec   = next((f for f in facts if f["Matched Fund Name"] == selected), None)
    if not rec or not rec.get("Category"):
        st.error(f"❌ Could not find category for '{selected}'.")
        return
    category = rec["Category"].strip()

    # --- 2) Load PPTX template
    try:
        prs = Presentation("assets/template.pptx")
    except Exception as e:
        st.error(f"❌ Could not load template: {e}")
        return
    slide1 = prs.slides[0]

    # --- 3) Titles on slides 1–4
    def set_title(slide, text):
        if slide.shapes.title:
            tx = slide.shapes.title.text_frame.paragraphs[0]
        else:
            left, top = Inches(0.5), Inches(0.3)
            box = slide.shapes.add_textbox(left, top, prs.slide_width - Inches(1), Inches(0.7))
            tx = box.text_frame.paragraphs[0]
        tx.text = text
        r = tx.runs[0]
        r.font.name  = "Helvetica"
        r.font.size  = Pt(20)
        r.font.bold  = True
        r.font.color.rgb = RGBColor(0x21, 0x2B, 0x58)
        tx.alignment = PP_ALIGN.LEFT

    set_title(slide1, "Investment Watchlist")
    if len(prs.slides) > 1:
        set_title(prs.slides[1], f"{category} – Expense & Return")
    if len(prs.slides) > 2:
        set_title(prs.slides[2], f"{category} – Risk Adjusted Statistics")
    if len(prs.slides) > 3:
        set_title(prs.slides[3], f"{category} – Qualitative Factors")

    # --- 4) Subheader on Slide 1
    sub = slide1.shapes.add_textbox(Inches(0.5), Inches(1.3), prs.slide_width - Inches(1), Inches(0.4))
    sf  = sub.text_frame.paragraphs[0]
    sf.text = selected
    run = sf.runs[0]
    run.font.name      = "Cambria"
    run.font.size      = Pt(12)
    run.font.bold      = True
    run.font.underline = True
    run.font.color.rgb = RGBColor(0, 0, 0)
    sf.alignment       = PP_ALIGN.LEFT

    # --- 5) Prepare bullets from Step 16 logic
    perf_data = st.session_state.get("fund_performance_data", [])
    item      = next((x for x in perf_data if x["Fund Scorecard Name"] == selected), {})

    template = st.session_state.get("bullet_point_templates", [""])[0]
    b1 = template
    for fld, val in item.items():
        b1 = b1.replace(f"[{fld}]", str(val))

    # --- IPS status: get from icon table for consistency
    ips_icon_table = st.session_state.get("ips_icon_table")
    ips_status = None
    if ips_icon_table is not None and not ips_icon_table.empty:
        row = ips_icon_table[ips_icon_table["Fund Name"] == selected]
        ips_status = row.iloc[0]["IPS Watch Status"] if not row.empty else None

    # --- Bullet 2: Watch status and return comparison
    if ips_status == "NW":
        b2 = "This fund is **not on watch**."
        b3 = ""
    else:
        status_label = (
            "Formal Watch" if ips_status == "FW" else
            "Informal Watch" if ips_status == "IW" else
            ips_status or "on watch"
        )
        three   = float(item.get("3Yr") or 0)
        bench3  = float(item.get("Bench 3Yr") or 0)
        five    = float(item.get("5Yr") or 0)
        bench5  = float(item.get("Bench 5Yr") or 0)
        bps3 = round((three - bench3) * 100, 1)
        bps5 = round((five  - bench5) * 100, 1)

        peer = st.session_state.get("step14_peer_rank_table", [])
        raw3 = next((r.get("Sharpe Ratio Rank 3Yr") for r in peer if r.get("Fund Name") == selected), None)
        raw5 = next((r.get("Sharpe Ratio Rank 5Yr") for r in peer if r.get("Fund Name") == selected), None)
        try:
            pos3 = "top" if raw3 and int(raw3) <= 50 else "bottom"
        except:
            pos3 = "bottom"
        try:
            pos5 = "top" if raw5 and int(raw5) <= 50 else "bottom"
        except:
            pos5 = "bottom"

        b2 = (
            f"The fund is now on **{status_label}**. Its 3‑year return trails the benchmark by "
            f"{bps3} bps ({three:.2f}% vs. {bench3:.2f}%) and its 5‑year return trails by "
            f"{bps5} bps ({five:.2f}% vs. {bench5:.2f}%). Its 3‑Yr Sharpe ranks in the {pos3} half of peers "
            f"and its 5‑Yr Sharpe ranks in the {pos5} half."
        )
        b3 = "**Action:** Consider replacing this fund." if ips_status == "FW" else ""

    # --- 6) Build Slide 1 Table (category, time, $plan assets, 1-11, status)
    row = {"Category": category, "Time Period": st.session_state.get("report_date", ""), "Plan Assets": "$"}
    # Use icon table to populate 1-11 & IPS status, match Slide 1 table
    if row and ips_icon_table is not None and not ips_icon_table.empty:
        icon_row = ips_icon_table[ips_icon_table["Fund Name"] == selected]
        if not icon_row.empty:
            icon_row = icon_row.iloc[0]
            for i in range(1, 12):
                row[str(i)] = icon_row.get(f"IPS Investment Criteria {i}", "")
            row["IPS Status"] = icon_row.get("IPS Watch Status", "")
        else:
            for i in range(1, 12):
                row[str(i)] = ""
            row["IPS Status"] = ""
    else:
        for i in range(1, 12):
            row[str(i)] = ""
        row["IPS Status"] = ""

    # --- Table setup
    slide_w = prs.slide_width
    table_w = Inches(9)
    left    = (slide_w - table_w) // 2
    top     = Inches(2.0)
    height  = Inches(0.6)
    cols    = 15
    col_w   = [1.2, 1.2, 1.2] + [0.4]*11 + [1.0]
    tbl     = slide1.shapes.add_table(2, cols, left, top, table_w, height).table

    # Column widths
    for idx, w in enumerate(col_w):
        tbl.columns[idx].width = Inches(w)

    # Headers/vals
    headers = ["Category", "Time Period", "Plan Assets"] + [str(i) for i in range(1, 12)] + ["IPS Status"]
    vals = [row["Category"], row["Time Period"], row["Plan Assets"]] + [row[str(i)] for i in range(1, 12)] + [row["IPS Status"]]

    # Helper: set border
    def _set_border(cell, color=RGBColor(0,0,0)):
        tc   = cell._tc
        tcPr = tc.get_or_add_tcPr()
        hex_val = "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])
        for ln_tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            ln   = OxmlElement(ln_tag)
            lnPr = OxmlElement("a:solidFill")
            srgb = OxmlElement("a:srgbClr"); srgb.set("val", hex_val)
            lnPr.append(srgb); ln.append(lnPr); tcPr.append(ln)

    # Header row formatting
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        _set_border(cell)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tfc = cell.text_frame; tfc.clear()
        tfc.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE; tfc.margin_top = tfc.margin_bottom = 0
        p = tfc.paragraphs[0]; r = p.add_run(); r.text = h
        r.font.name = "Cambria"; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

    # Data row formatting (icons for ✔/✖)
    for c, val in enumerate(vals):
        cell = tbl.cell(1, c)
        _set_border(cell)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tfc = cell.text_frame; tfc.clear()
        tfc.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE; tfc.margin_top = tfc.margin_bottom = 0
        p = tfc.paragraphs[0]; p.font.name = "Cambria"; p.font.size = Pt(12); p.alignment = PP_ALIGN.CENTER

        if c < 14:
            if val == "✔":
                p.text = "✔"; p.runs[0].font.color.rgb = RGBColor(0,176,80)
            elif val == "✗":
                p.text = "✖"; p.runs[0].font.color.rgb = RGBColor(192,0,0)
            else:
                p.text = str(val)
        else:
            # IPS Status badge
            badge, color = "", RGBColor(0,176,80)
            if vals[c] == "IW":
                badge, color = "IW", RGBColor(255,165,0)
            elif vals[c] == "FW":
                badge, color = "FW", RGBColor(192,0,0)
            elif vals[c] == "NW":
                badge, color = "NW", RGBColor(0,176,80)
            else:
                badge, color = str(vals[c]), RGBColor(200,200,200)
            badge_width = Inches(0.6)
            badge_height = Inches(0.4)
            cell_left  = left + sum(Inches(w) for w in col_w[:c])
            cell_w     = Inches(col_w[c])
            bx = cell_left + (cell_w - badge_width) / 2
            by = top + (height - badge_height) / 2 + Inches(0.08)
            shp = slide1.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, badge_width, badge_height)
            shp.fill.solid(); shp.fill.fore_color.rgb = color
            shp.line.color.rgb = RGBColor(255, 255, 255)
            tf2 = shp.text_frame; tf2.clear()
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            p2.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            r2 = p2.add_run(); r2.text = badge
            r2.font.name = "Cambria"; r2.font.size = Pt(10); r2.font.bold = True; r2.font.color.rgb = RGBColor(255, 255, 255)

    # --- 7) Bullet textbox BELOW the table
    bullet_gap  = Inches(0.1)
    bullet_left = left
    bullet_top  = top + height + bullet_gap
    bullet_w    = prs.slide_width - Inches(1)
    bullet_h    = Inches(2.5)

    tb = slide1.shapes.add_textbox(bullet_left, bullet_top, bullet_w, bullet_h)
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    for text in [b1, b2] + ([b3] if b3 else []):
        p = tf.add_paragraph(); p.text = f"• {text}"
        p.font.name    = "Cambria"; p.font.size = Pt(11); p.font.bold = text.startswith("**Action")
        p.alignment    = PP_ALIGN.LEFT; p.line_spacing = 2.0

    # === Slide 2: Expense & Return ===
    if len(prs.slides) > 1:
        slide2 = prs.slides[1]
        
        # Table 1: Net Expense Ratio
        perf_item = next(p for p in perf_data if p["Fund Scorecard Name"] == selected)
        inv_mgr   = f"{selected} ({perf_item['Ticker']})"
        net_exp   = perf_item.get("Net Expense Ratio", "")
        if net_exp and not str(net_exp).endswith("%"):
            net_exp = f"{net_exp}%"
        df1 = pd.DataFrame([{"Investment Manager": inv_mgr, "Net Expense Ratio": net_exp}])
    
        # Table 2: QTD / 1Yr / 3Yr / 5Yr / 10Yr
        date_label = st.session_state.get("report_date", "QTD")
        def _pct(v): return f"{v}%" if v and not str(v).endswith("%") else (v or "")
        df2 = pd.DataFrame([{
            "Investment Manager": inv_mgr,
            date_label:           _pct(perf_item.get("QTD")),
            "1 Year":             _pct(perf_item.get("1Yr")),
            "3 Year":             _pct(perf_item.get("3Yr")),
            "5 Year":             _pct(perf_item.get("5Yr")),
            "10 Year":            _pct(perf_item.get("10Yr")),
        }])
    
        # Table 3: Last 10 Calendar‑Year Returns
        fund_cy  = st.session_state.get("step8_returns", [])
        bench_cy = st.session_state.get("benchmark_calendar_year_returns", [])
        fund_rec = next((r for r in fund_cy if r.get("Name") == selected), {})
        bench_rec= next((r for r in bench_cy if r.get("Name") == selected or r.get("Ticker") == fund_rec.get("Ticker")), {})
        years    = sorted([c for c in fund_rec.keys() if re.match(r"20\d{2}", c)], reverse=True)[:10]
        rows3 = [
            {"Investment Manager": inv_mgr, **{y: fund_rec.get(y, "") for y in years}},
            {"Investment Manager": f"{bench_rec.get('Name','Benchmark')} ({bench_rec.get('Ticker','')})",
             **{y: bench_rec.get(y, "") for y in years}}
        ]
        df3 = pd.DataFrame(rows3, columns=["Investment Manager"] + years)
    
        def draw_table(slide, df, left, top, width, height, col_widths):
            tbl = slide.shapes.add_table(len(df)+1, len(df.columns),
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(height)).table
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Inches(w)
            # header
            for c, name in enumerate(df.columns):
                cell = tbl.cell(0, c); cell.text = name
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(33,43,88)  # Dark blue header
                p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                run = p.runs[0]; run.font.name="Cambria"; run.font.size=Pt(12); run.font.bold=True
                run.font.color.rgb = RGBColor(255,255,255)  # White text for the header
            # body
            for r in range(len(df)):
                for c in range(len(df.columns)):
                    cell = tbl.cell(r+1, c); cell.text = str(df.iat[r, c])
                    if r % 2 == 0:
                        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(240,245,255)
                    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]; run.font.name="Cambria"; run.font.size=Pt(12); run.font.color.rgb = RGBColor(0,0,0)
                    # Dark blue for the first column
                    if c == 0:
                        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(33,43,88)
                        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                        run = p.runs[0]; run.font.color.rgb = RGBColor(255,255,255)  # White text for the first column
    
        # Get slide width and margins
        sw      = prs.slide_width / Inches(1)  # Slide width in inches
        lm, rm, gap = 0.5, 0.5, 0.2  # Left margin, right margin, and gap between tables
        usable = sw - lm - rm - gap  # Usable space in width
        half   = usable / 2  # Half of the usable width for splitting the table
        ty, hgh = 1.0, 1.2  # Starting top position (ty) and height (hgh) for the first table
    
        # Adjust the vertical position for Slide 2 Table 3 (increased by 3.0 to push it lower)
        by, bh  = ty + hgh + 3.0, 2.0  # Increased the vertical distance from the previous table (was 0.3)
    
        # Draw Table 1 (Net Expense Ratio) - Positioned at the top of the slide
        draw_table(slide2, df1, lm, ty, half, hgh, [2.0, half-2.0])
    
        # Draw Table 2 (QTD / 1Yr / 3Yr / 5Yr / 10Yr) - Positioned to the right of Table 1
        draw_table(slide2, df2, lm+half+gap, ty, half, hgh, [2.0] + [(half-2.0)/5]*5)
    
        # Prepare and draw Table 3 (Last 10 Calendar-Year Returns) - Positioned further down
        first_w = 2.5  # Starting width for the first column in Table 3
        extras  = len(df3.columns) - 1  # Number of extra columns to be evenly distributed
        cw3     = [first_w] + ([(usable-first_w)/extras]*extras if extras > 0 else [])  # Calculate column widths for Table 3
        draw_table(slide2, df3, lm, by, usable, bh, cw3)  # Draw Table 3 at the new adjusted position


    # 9) Slide 3: Risk-Adjusted Statistics
    if len(prs.slides) > 2:
        slide3 = prs.slides[2]
        inv_mgr_label = selected  # will append ticker below

        # === Table 1: 3-Yr & 5-Yr MPT Stats ===
        mpt3 = st.session_state.get("step9_mpt_stats", [])
        stats3 = next((r for r in mpt3 if r["Fund Name"] == selected), {})
        mpt5 = st.session_state.get("step10_mpt_stats", [])
        stats5 = next((r for r in mpt5 if r["Fund Name"] == selected), {})

        ticker = stats3.get("Ticker", stats5.get("Ticker", ""))
        inv_mgr = f"{selected} ({ticker})"

        row1 = {
            "Investment Manager":        inv_mgr,
            "3 Year Alpha":              stats3.get("3 Year Alpha", ""),
            "5 Year Alpha":              stats5.get("5 Year Alpha", ""),
            "3 Year Beta":               stats3.get("3 Year Beta", ""),
            "5 Year Beta":               stats5.get("5 Year Beta", ""),
            "3 Year Upside Capture":     stats3.get("3 Year Upside Capture", ""),
            "3 Year Downside Capture":   stats3.get("3 Year Downside Capture", ""),
            "5 Year Upside Capture":     stats5.get("5 Year Upside Capture", ""),
            "5 Year Downside Capture":   stats5.get("5 Year Downside Capture", "")
        }
        df3_1 = pd.DataFrame([row1])
        cols1 = list(df3_1.columns)
        widths1 = [2.5] + [(prs.slide_width.inches - 2.5) / (len(cols1)-1)] * (len(cols1)-1)

        # === Table 2: Ratios & Peer Ranks ===
        risk_table = st.session_state.get("step13_risk_adjusted_table", [])
        peer_table = st.session_state.get("step14_peer_rank_table", [])
        risk_rec = next((r for r in risk_table if r["Fund Name"] == selected), {})
        peer_rec = next((r for r in peer_table if r["Fund Name"] == selected), {})

        def frac(metric, period):
            return f"{risk_rec.get(f'{metric} {period}', '')} / {peer_rec.get(f'{metric} {period}', '')}"

        row2 = {
            "Investment Manager": inv_mgr,
            "3Yr Sharpe / Peer %":      frac("Sharpe Ratio", "3Yr"),
            "5Yr Sharpe / Peer %":      frac("Sharpe Ratio", "5Yr"),
            "3Yr Sortino / Peer %":     frac("Sortino Ratio", "3Yr"),
            "5Yr Sortino / Peer %":     frac("Sortino Ratio", "5Yr"),
            "3Yr Info Ratio / Peer %":  frac("Information Ratio", "3Yr"),
            "5Yr Info Ratio / Peer %":  frac("Information Ratio", "5Yr"),
        }
        df3_2 = pd.DataFrame([row2])
        cols2 = list(df3_2.columns)
        widths2 = [2.5] + [(prs.slide_width.inches - 2.5) / (len(cols2)-1)] * (len(cols2)-1)

        # helper for drawing
        def draw_table(slide, df, left, top, width, height, col_widths):
            tbl = slide.shapes.add_table(len(df)+1, len(df.columns),
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(height)).table
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Inches(w)
            # header
            for c, name in enumerate(df.columns):
                cell = tbl.cell(0, c)
                cell.text = name
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(33,43,88)
                p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                run = p.runs[0]; run.font.name="Cambria"; run.font.size=Pt(12); run.font.bold=True
                run.font.color.rgb = RGBColor(255,255,255)
            # body
            for r in range(len(df)):
                for c in range(len(df.columns)):
                    cell = tbl.cell(r+1, c)
                    cell.text = str(df.iat[r, c])
                    if r % 2 == 0:
                        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(240,245,255)
                    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                    run = p.runs[0]; run.font.name="Cambria"; run.font.size=Pt(12); run.font.color.rgb = RGBColor(0,0,0)

        # calculate positions
        sw = prs.slide_width.inches
        left1, top1, h1 = 0.5, 1.0, 1.0
        left2, top2, h2 = 0.5, 2.3, 1.0
        total_w = sw - 1.0

        draw_table(slide3, df3_1, left1, top1, total_w, h1, widths1)
        draw_table(slide3, df3_2, left2, top2, total_w, h2, widths2)

    # 4) Slide 4: Qualitative Factors Tables side by side
    if len(prs.slides) > 3:
        slide4 = prs.slides[3]
        # Prepare data for Table 1
        blocks = st.session_state.get("fund_blocks", [])
        block = next((b for b in blocks if b["Fund Name"] == selected), {})
        raw_tenure = next((m["Info"] for m in block.get("Metrics", []) if m["Metric"] == "Manager Tenure"), "")
        m = re.search(r"(\d+(\.\d+)?)", raw_tenure)
        tenure = f"{m.group(1)} years" if m else raw_tenure
        perf_data = st.session_state.get("fund_performance_data", [])
        perf_item = next((p for p in perf_data if p.get("Fund Scorecard Name") == selected), {})
        ticker = perf_item.get("Ticker", "")
        inv_mgr = f"{selected} ({ticker})"
        df4_1 = pd.DataFrame([{"Investment Manager": inv_mgr, "Manager Tenure": tenure}])
        # Prepare data for Table 2
        facts = st.session_state.get("fund_factsheets_data", [])
        fs_rec = next((f for f in facts if f["Matched Fund Name"] == selected), {})
        assets = fs_rec.get("Net Assets", "")
        avg_cap = fs_rec.get("Avg. Market Cap", "")
        df4_2 = pd.DataFrame([{
            "Investment Manager": inv_mgr,
            "Assets Under Management": assets,
            "Average Market Capitalization": avg_cap
        }])

        # Helper to draw a table with header & first column dark blue and alternating stripes
        def draw_table(slide, df, left, top, width, height, col_widths):
            tbl = slide.shapes.add_table(len(df)+1, len(df.columns),
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(height)).table
            # Set column widths
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Inches(w)
            # Header row: all dark blue, white text
            for c, name in enumerate(df.columns):
                cell = tbl.cell(0, c)
                cell.text = name
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(33,43,88)
                p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.runs[0]; r.font.name="Cambria"; r.font.size=Pt(12); r.font.bold=True; r.font.color.rgb=RGBColor(255,255,255)
            # Body rows
            for r in range(len(df)):
                for c in range(len(df.columns)):
                    cell = tbl.cell(r+1, c)
                    cell.text = str(df.iat[r, c])
                    cell.fill.solid()
                    # First column always dark blue
                    if c == 0:
                        cell.fill.fore_color.rgb = RGBColor(33,43,88)
                        text_color = RGBColor(255,255,255)
                    else:
                        # alternating stripe on other cols
                        if r % 2 == 0:
                            cell.fill.fore_color.rgb = RGBColor(240,245,255)
                        else:
                            cell.fill.fore_color.rgb = RGBColor(255,255,255)
                        text_color = RGBColor(0,0,0)
                    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                    rfont = p.runs[0]; rfont.font.name="Cambria"; rfont.font.size=Pt(12); rfont.font.color.rgb=text_color

        # Calculate positions & widths
        sw = prs.slide_width.inches
        margin = 0.5
        gap = 0.2
        usable = sw - 2*margin - gap
        half = usable / 2
        left1 = margin
        left2 = margin + half + gap
        top = 1.3  # lowered position (move tables down)
        height = 0.8
        # Column widths: Table1 (manager wide, tenure small), Table2 (manager fixed, others larger)
        cw1 = [half - 1.5, 1.5]
        cw2 = [2.0, (half - 2.0)*0.3, (half - 2.0)*0.7]  # allocate more width to Average Market Capitalization

        draw_table(slide4, df4_1, left1, top, half, height, cw1)
        draw_table(slide4, df4_2, left2, top, half, height, cw2)

    # 5) Download
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    st.download_button(
        label="Download PPTX",
        data=buf,
        file_name=f"{selected.replace(' ','_')}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key=f"download_ppt_{selected}"
    )

#─────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
# === Main App ===
def run():
    import re
    st.title("Writeup")
    uploaded = st.file_uploader("Upload MPI PDF", type="pdf")
    if not uploaded:
        return
   #──────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────────
    
    with pdfplumber.open(uploaded) as pdf:
        # Step 1
        with st.expander("Step 1: Details", expanded=False):
            first = pdf.pages[0].extract_text() or ""
            process_page1(first)

        # Step 2
        with st.expander("Step 2: Table of Contents", expanded=False):
            toc_text = "".join((pdf.pages[i].extract_text() or "") for i in range(min(3, len(pdf.pages))))
            process_toc(toc_text)

        # --- COMBINED STEPS 3, 4, 5 ---
        with st.expander("Step 3: Scorecard + IPS + Fund Type", expanded=True):
            sp = st.session_state.get('scorecard_page')
            tot = st.session_state.get('total_options')
            pp = st.session_state.get('performance_page')
            factsheets_page = st.session_state.get('factsheets_page')
            if sp and tot is not None and pp:
                step3_5_6_scorecard_and_ips(pdf, sp, pp, factsheets_page, tot)
            else:
                st.error("Missing scorecard, performance page, or total options")

        
        # Step 6
        with st.expander("Step 6: Fund Factsheets", expanded=True):
            names = [b['Fund Name'] for b in st.session_state.get('fund_blocks', [])]
            step6_process_factsheets(pdf, names)

        # Step 7
        with st.expander("Step 7: Annualized Returns", expanded=False):
            step7_extract_returns(pdf)

        # ── Data Prep for Bullet Points ───────────────────────────────────────────────
        report_date = st.session_state.get("report_date", "")
        m = re.match(r"(\d)(?:st|nd|rd|th)\s+QTR,\s*(\d{4})", report_date)
        quarter = m.group(1) if m else ""
        year    = m.group(2) if m else ""
        
        # make sure we have a list to iterate
        for itm in st.session_state.get("fund_performance_data", []):
            # force numeric defaults
            qtd       = float(itm.get("QTD") or 0)
            bench_qtd = float(itm.get("Bench QTD") or 0)
        
            # direction, quarter/year
            itm["Perf Direction"] = "overperformed" if qtd >= bench_qtd else "underperformed"
            itm["Quarter"]        = quarter
            itm["Year"]           = year
        
            # basis‑points difference
            diff_bps = round((qtd - bench_qtd) * 100, 1)
            itm["QTD_bps_diff"] = str(diff_bps)
        
            # percent strings
            fund_pct  = f"{qtd:.2f}%"
            bench_pct = f"{bench_qtd:.2f}%"
            itm["QTD_pct_diff"] = f"{(qtd - bench_qtd):.2f}%"
            itm["QTD_vs"]       = f"{fund_pct} vs. {bench_pct}"
        
        # Initialize your template exactly once
        if "bullet_point_templates" not in st.session_state:
            st.session_state["bullet_point_templates"] = [
                "[Fund Scorecard Name] [Perf Direction] its benchmark in Q[Quarter], "
                "[Year] by [QTD_bps_diff] bps ([QTD_vs])."
            ]

        # ───────────────────────────────────────────────────────────────────────────────
        
        # Step 8: Calendar Year Section
        with st.expander("Step 8: Calendar Year Returns", expanded=False):
            step8_calendar_returns(pdf)

        # Step 9: Match Tickers
        with st.expander("Step 9: Risk Analysis (3Yr)", expanded=False):
            step9_risk_analysis_3yr(pdf)

        # Step 10: Match Tickers
        with st.expander("Step 10: Risk Analysis (5Yr)", expanded=False):
            step10_risk_analysis_5yr(pdf)

        # Step 11: MPT Statistics Summary
        with st.expander("Step 11: MPT Statistics Summary", expanded=False):
            step11_create_summary()
            
        # Step 12: Find Factsheet Sub‑Headings
        with st.expander("Step 12: Fund Facts ", expanded=False):
            step12_process_fund_facts(pdf)

        # Step 13: Risk Adjusted Returns
        with st.expander("Step 13: Risk-Adjusted Returns", expanded=False):
            step13_process_risk_adjusted_returns(pdf)

        # Step 14: Peer Risk-Adjusted Return Rank
        with st.expander("Step 14: Peer Risk-Adjusted Return Rank", expanded=False):
            step14_extract_peer_risk_adjusted_return_rank(pdf)
        
        # Step 15: View Single Fund Details
        with st.expander("Step 15: Single Fund Details", expanded=False):
            step15_display_selected_fund()
                    
        # Step 16: Bullet Points
        with st.expander("Step 16: Bullet Points", expanded=False):
            step16_bullet_points()
            
        with st.expander("Step 17: Export to PowerPoint", expanded=False):
            step17_export_to_ppt_headings()


if __name__ == "__main__":
    run()
